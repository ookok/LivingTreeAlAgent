"""
HuashuDesignTool - 华数设计集成工具

提供流程图绘制、总平面图绘制和 PPT 生成等功能。

功能模块：
1. 流程图绘制 - 支持多种流程图类型（流程、时序、泳道、思维导图等）
2. 总平面图绘制 - 支持建筑、景观、规划等场景的平面图绘制
3. PPT 生成 - 根据内容自动生成演示文稿

Author: LivingTreeAI Agent
Date: 2026-04-30
"""

import os
import json
from typing import Dict, Any, Optional, List, Union
from dataclasses import dataclass, field
from loguru import logger
from enum import Enum

from business.tools.base_tool import BaseTool


class DiagramType(Enum):
    """流程图类型枚举"""
    FLOWCHART = "flowchart"           # 标准流程图
    SEQUENCE = "sequence"             # 时序图
    SWIMLANE = "swimlane"             # 泳道图
    MINDMAP = "mindmap"               # 思维导图
    ER = "er"                         # ER 图
    CLASS = "class"                   # 类图
    NETWORK = "network"               # 网络图


class LayoutType(Enum):
    """布局类型枚举"""
    TOP_DOWN = "top_down"             # 自上而下
    LEFT_RIGHT = "left_right"         # 自左向右
    RADIAL = "radial"                 # 放射状（思维导图）
    ORGANIC = "organic"               # 有机布局


class ElementType(Enum):
    """流程图元素类型"""
    START = "start"                   # 开始节点
    END = "end"                       # 结束节点
    PROCESS = "process"               # 处理节点
    DECISION = "decision"             # 决策节点
    INPUT_OUTPUT = "input_output"     # 输入输出节点
    CONNECTOR = "connector"           # 连接节点
    SUBPROCESS = "subprocess"         # 子流程
    LOOP = "loop"                     # 循环节点


class ColorTheme(Enum):
    """配色主题"""
    DEFAULT = "default"               # 默认主题
    BLUE = "blue"                     # 蓝色主题
    GREEN = "green"                   # 绿色主题
    ORANGE = "orange"                 # 橙色主题
    PURPLE = "purple"                 # 紫色主题
    GRAPHITE = "graphite"             # 石墨主题


@dataclass
class FlowNode:
    """流程图节点"""
    id: str
    type: str
    label: str
    x: Optional[int] = None
    y: Optional[int] = None
    width: int = 120
    height: int = 60
    color: Optional[str] = None
    style: Optional[Dict[str, Any]] = None


@dataclass
class FlowEdge:
    """流程图边"""
    id: str
    source: str
    target: str
    label: Optional[str] = None
    style: Optional[Dict[str, Any]] = None


@dataclass
class DiagramConfig:
    """图表配置"""
    title: str = ""
    type: str = "flowchart"
    layout: str = "top_down"
    theme: str = "default"
    width: int = 800
    height: int = 600
    show_grid: bool = False
    show_labels: bool = True


@dataclass
class FloorPlanElement:
    """平面图元素"""
    id: str
    type: str
    label: str
    x: int = 0
    y: int = 0
    width: int = 100
    height: int = 100
    rotation: float = 0.0
    color: str = "#3498db"
    properties: Optional[Dict[str, Any]] = None


@dataclass
class PPTSlide:
    """PPT 幻灯片"""
    title: str
    content: str
    layout: str = "title_content"  # title_only, title_content, content_only, comparison, blank
    image_path: Optional[str] = None
    transition: str = "fade"


@dataclass
class PPTConfig:
    """PPT 配置"""
    title: str = "演示文稿"
    author: str = "LivingTreeAI"
    theme: str = "default"
    slides: List[PPTSlide] = field(default_factory=list)


class HuashuDesignTool(BaseTool):
    """
    华数设计集成工具
    
    功能：
    1. 流程图绘制 - 支持多种流程图类型
    2. 总平面图绘制 - 支持建筑、景观、规划等场景
    3. PPT 生成 - 根据内容自动生成演示文稿
    
    使用方式：
        tool = HuashuDesignTool()
        
        # 绘制流程图
        result = await tool.execute(
            action="create_flowchart",
            nodes=[...],
            edges=[...],
            config={"title": "工作流程图"}
        )
        
        # 绘制平面图
        result = await tool.execute(
            action="create_floorplan",
            elements=[...],
            config={"title": "总平面图"}
        )
        
        # 生成 PPT
        result = await tool.execute(
            action="generate_ppt",
            content="演示内容描述...",
            output_path="/path/to/output.pptx"
        )
    """
    
    def __init__(self):
        super().__init__()
        self._logger = logger.bind(component="HuashuDesignTool")
        self._output_dir = self._get_output_dir()
    
    @property
    def name(self) -> str:
        return "huashu_design"
    
    @property
    def description(self) -> str:
        return "华数设计集成工具，提供流程图绘制（流程/时序/泳道/思维导图等）、总平面图绘制（建筑/景观/规划）和 PPT 自动生成功能"
    
    @property
    def category(self) -> str:
        return "design"
    
    @property
    def node_type(self) -> str:
        return "deterministic"
    
    def _get_output_dir(self) -> str:
        """获取输出目录"""
        output_dir = os.path.join(os.path.expanduser("~"), "LivingTreeAI", "design_output")
        os.makedirs(output_dir, exist_ok=True)
        return output_dir
    
    def _generate_id(self) -> str:
        """生成唯一 ID"""
        import uuid
        return str(uuid.uuid4())[:8]
    
    async def execute(self, action: str, **kwargs) -> Dict[str, Any]:
        """
        执行设计操作
        
        Args:
            action: 操作类型（create_flowchart/create_floorplan/generate_ppt/render）
            **kwargs: 操作参数
            
        Returns:
            操作结果
        """
        try:
            if action == "create_flowchart":
                return await self._create_flowchart(**kwargs)
            elif action == "create_floorplan":
                return await self._create_floorplan(**kwargs)
            elif action == "generate_ppt":
                return await self._generate_ppt(**kwargs)
            elif action == "render":
                return await self._render_diagram(**kwargs)
            elif action == "list_templates":
                return await self._list_templates(**kwargs)
            else:
                return {
                    "success": False,
                    "message": f"未知操作类型: {action}",
                    "data": None
                }
        except Exception as e:
            self._logger.error(f"执行设计操作失败: {e}")
            return {
                "success": False,
                "message": str(e),
                "data": None
            }
    
    async def _create_flowchart(
        self,
        nodes: Optional[List[Dict[str, Any]]] = None,
        edges: Optional[List[Dict[str, Any]]] = None,
        config: Optional[Dict[str, Any]] = None,
        **kwargs
    ) -> Dict[str, Any]:
        """
        创建流程图
        
        Args:
            nodes: 节点列表，每个节点包含 id, type, label, x, y, width, height, color
            edges: 边列表，每个边包含 id, source, target, label
            config: 配置参数（title, type, layout, theme, width, height）
            
        Returns:
            流程图数据
        """
        if not nodes or not edges:
            return {
                "success": False,
                "message": "节点和边数据不能为空",
                "data": None
            }
        
        # 验证节点类型
        valid_types = [e.value for e in ElementType]
        for node in nodes:
            if node.get("type") not in valid_types:
                return {
                    "success": False,
                    "message": f"无效的节点类型: {node.get('type')}，有效值: {valid_types}",
                    "data": None
                }
        
        # 构建流程图数据
        diagram_data = {
            "type": config.get("type", "flowchart"),
            "title": config.get("title", "流程图"),
            "layout": config.get("layout", "top_down"),
            "theme": config.get("theme", "default"),
            "width": config.get("width", 800),
            "height": config.get("height", 600),
            "nodes": nodes,
            "edges": edges,
            "metadata": {
                "created_at": self._get_timestamp(),
                "version": "1.0"
            }
        }
        
        # 自动布局（如果没有指定坐标）
        if not all(node.get("x") is not None and node.get("y") is not None for node in nodes):
            diagram_data["nodes"] = self._auto_layout(nodes, edges, diagram_data["layout"])
        
        self._logger.info(f"创建流程图: {diagram_data['title']}, {len(nodes)} 个节点, {len(edges)} 条边")
        
        return {
            "success": True,
            "message": "流程图创建成功",
            "data": diagram_data,
            "diagram_id": self._generate_id()
        }
    
    def _auto_layout(self, nodes: List[Dict[str, Any]], edges: List[Dict[str, Any]], 
                     layout_type: str = "top_down") -> List[Dict[str, Any]]:
        """
        自动布局节点
        
        Args:
            nodes: 节点列表
            edges: 边列表
            layout_type: 布局类型
            
        Returns:
            带有坐标的节点列表
        """
        node_width = 120
        node_height = 60
        spacing_x = 150
        spacing_y = 100
        start_x = 100
        start_y = 80
        
        # 构建节点位置映射
        node_positions = {}
        current_x = start_x
        current_y = start_y
        
        if layout_type == "top_down":
            # 自上而下布局
            for i, node in enumerate(nodes):
                node_id = node["id"]
                node_positions[node_id] = {
                    "x": current_x,
                    "y": current_y + i * (node_height + spacing_y)
                }
        
        elif layout_type == "left_right":
            # 自左向右布局
            for i, node in enumerate(nodes):
                node_id = node["id"]
                node_positions[node_id] = {
                    "x": current_x + i * (node_width + spacing_x),
                    "y": current_y
                }
        
        else:
            # 默认自上而下
            for i, node in enumerate(nodes):
                node_id = node["id"]
                node_positions[node_id] = {
                    "x": current_x,
                    "y": current_y + i * (node_height + spacing_y)
                }
        
        # 更新节点坐标
        result_nodes = []
        for node in nodes:
            pos = node_positions.get(node["id"], {"x": current_x, "y": current_y})
            result_nodes.append({
                **node,
                "x": node.get("x", pos["x"]),
                "y": node.get("y", pos["y"]),
                "width": node.get("width", node_width),
                "height": node.get("height", node_height)
            })
        
        return result_nodes
    
    async def _create_floorplan(
        self,
        elements: Optional[List[Dict[str, Any]]] = None,
        config: Optional[Dict[str, Any]] = None,
        **kwargs
    ) -> Dict[str, Any]:
        """
        创建总平面图
        
        Args:
            elements: 平面图元素列表
            config: 配置参数
            
        Returns:
            平面图数据
        """
        if not elements:
            return {
                "success": False,
                "message": "平面图元素不能为空",
                "data": None
            }
        
        # 构建平面图数据
        floorplan_data = {
            "type": "floorplan",
            "title": config.get("title", "总平面图"),
            "width": config.get("width", 1200),
            "height": config.get("height", 800),
            "scale": config.get("scale", 1.0),
            "units": config.get("units", "米"),
            "elements": elements,
            "metadata": {
                "created_at": self._get_timestamp(),
                "version": "1.0"
            }
        }
        
        self._logger.info(f"创建总平面图: {floorplan_data['title']}, {len(elements)} 个元素")
        
        return {
            "success": True,
            "message": "总平面图创建成功",
            "data": floorplan_data,
            "plan_id": self._generate_id()
        }
    
    async def _generate_ppt(
        self,
        content: Optional[str] = None,
        slides: Optional[List[Dict[str, Any]]] = None,
        output_path: Optional[str] = None,
        config: Optional[Dict[str, Any]] = None,
        **kwargs
    ) -> Dict[str, Any]:
        """
        生成 PPT 演示文稿
        
        Args:
            content: 内容描述（用于自动生成幻灯片）
            slides: 幻灯片列表（手动指定）
            output_path: 输出路径
            config: 配置参数
            
        Returns:
            PPT 生成结果
        """
        if not content and not slides:
            return {
                "success": False,
                "message": "内容描述或幻灯片列表不能为空",
                "data": None
            }
        
        # 确定输出路径
        if not output_path:
            timestamp = self._get_timestamp().replace(":", "-")
            output_path = os.path.join(self._output_dir, f"presentation_{timestamp}.pptx")
        
        # 确保输出目录存在
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        # 如果提供了内容描述，自动生成幻灯片
        if content and not slides:
            slides = self._generate_slides_from_content(content)
        
        # 构建 PPT 数据
        ppt_data = {
            "title": config.get("title", "演示文稿"),
            "author": config.get("author", "LivingTreeAI"),
            "theme": config.get("theme", "default"),
            "slides": slides,
            "metadata": {
                "created_at": self._get_timestamp(),
                "version": "1.0",
                "output_path": output_path
            }
        }
        
        # 生成 PPT 文件
        success = self._render_ppt(ppt_data, output_path)
        
        if success:
            self._logger.info(f"PPT 生成成功: {output_path}")
            return {
                "success": True,
                "message": "PPT 生成成功",
                "data": ppt_data,
                "output_path": output_path,
                "slide_count": len(slides)
            }
        else:
            return {
                "success": False,
                "message": "PPT 生成失败",
                "data": ppt_data
            }
    
    def _generate_slides_from_content(self, content: str) -> List[Dict[str, Any]]:
        """
        根据内容描述自动生成幻灯片
        
        Args:
            content: 内容描述
            
        Returns:
            幻灯片列表
        """
        slides = []
        
        # 简单的内容解析（可扩展为更复杂的 NLP 分析）
        paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
        
        if paragraphs:
            # 封面页
            slides.append({
                "title": paragraphs[0][:50] if len(paragraphs[0]) > 50 else paragraphs[0],
                "content": "",
                "layout": "title_only"
            })
            
            # 内容页
            for i, paragraph in enumerate(paragraphs[1:], 1):
                slides.append({
                    "title": f"第 {i} 节",
                    "content": paragraph,
                    "layout": "title_content"
                })
            
            # 结尾页
            slides.append({
                "title": "谢谢观看",
                "content": "",
                "layout": "title_only"
            })
        
        return slides
    
    def _render_ppt(self, ppt_data: Dict[str, Any], output_path: str) -> bool:
        """
        渲染 PPT 文件
        
        Args:
            ppt_data: PPT 数据
            output_path: 输出路径
            
        Returns:
            是否成功
        """
        try:
            # 尝试使用 python-pptx 库生成 PPT
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.dml.color import RGBColor
                
                prs = Presentation()
                
                # 设置标题
                prs.core_properties.title = ppt_data["title"]
                prs.core_properties.author = ppt_data["author"]
                
                for slide_data in ppt_data["slides"]:
                    # 根据布局类型选择幻灯片布局
                    if slide_data["layout"] == "title_only":
                        slide_layout = prs.slide_layouts[0]  # 标题页
                    elif slide_data["layout"] == "title_content":
                        slide_layout = prs.slide_layouts[1]  # 标题+内容
                    else:
                        slide_layout = prs.slide_layouts[1]
                    
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # 设置标题
                    title = slide.shapes.title
                    title.text = slide_data["title"]
                    
                    # 设置内容
                    if slide_data["content"]:
                        content = slide.placeholders[1]
                        content.text = slide_data["content"]
                
                prs.save(output_path)
                return True
                
            except ImportError:
                # 如果没有 python-pptx，生成 JSON 格式的 PPT 数据文件
                json_path = output_path.replace(".pptx", ".json")
                with open(json_path, 'w', encoding='utf-8') as f:
                    json.dump(ppt_data, f, ensure_ascii=False, indent=2)
                self._logger.warning("python-pptx 未安装，生成 JSON 格式数据")
                return True
                
        except Exception as e:
            self._logger.error(f"渲染 PPT 失败: {e}")
            return False
    
    async def _render_diagram(
        self,
        diagram_data: Optional[Dict[str, Any]] = None,
        output_path: Optional[str] = None,
        format: str = "svg",
        **kwargs
    ) -> Dict[str, Any]:
        """
        渲染图表为图片
        
        Args:
            diagram_data: 图表数据
            output_path: 输出路径
            format: 输出格式（svg, png, json）
            
        Returns:
            渲染结果
        """
        if not diagram_data:
            return {
                "success": False,
                "message": "图表数据不能为空",
                "data": None
            }
        
        # 确定输出路径
        if not output_path:
            timestamp = self._get_timestamp().replace(":", "-")
            output_path = os.path.join(self._output_dir, f"diagram_{timestamp}.{format}")
        
        # 确保输出目录存在
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        try:
            if format == "json":
                # 输出 JSON 格式
                with open(output_path, 'w', encoding='utf-8') as f:
                    json.dump(diagram_data, f, ensure_ascii=False, indent=2)
                
            elif format == "svg":
                # 生成 SVG
                svg_content = self._generate_svg(diagram_data)
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(svg_content)
            
            else:
                # 其他格式暂不支持，输出 JSON
                json_path = output_path.replace(f".{format}", ".json")
                with open(json_path, 'w', encoding='utf-8') as f:
                    json.dump(diagram_data, f, ensure_ascii=False, indent=2)
                output_path = json_path
            
            self._logger.info(f"图表渲染成功: {output_path}")
            
            return {
                "success": True,
                "message": "图表渲染成功",
                "data": diagram_data,
                "output_path": output_path,
                "format": format
            }
            
        except Exception as e:
            self._logger.error(f"渲染图表失败: {e}")
            return {
                "success": False,
                "message": str(e),
                "data": None
            }
    
    def _generate_svg(self, diagram_data: Dict[str, Any]) -> str:
        """
        生成 SVG 格式图表
        
        Args:
            diagram_data: 图表数据
            
        Returns:
            SVG 字符串
        """
        width = diagram_data.get("width", 800)
        height = diagram_data.get("height", 600)
        nodes = diagram_data.get("nodes", [])
        edges = diagram_data.get("edges", [])
        
        svg_parts = []
        svg_parts.append(f'<svg width="{width}" height="{height}" xmlns="http://www.w3.org/2000/svg">')
        
        # 绘制边
        for edge in edges:
            source_node = next((n for n in nodes if n["id"] == edge["source"]), None)
            target_node = next((n for n in nodes if n["id"] == edge["target"]), None)
            
            if source_node and target_node:
                sx = source_node.get("x", 0) + source_node.get("width", 120) / 2
                sy = source_node.get("y", 0) + source_node.get("height", 60) / 2
                tx = target_node.get("x", 0) + target_node.get("width", 120) / 2
                ty = target_node.get("y", 0) + target_node.get("height", 60) / 2
                
                svg_parts.append(f'<line x1="{sx}" y1="{sy}" x2="{tx}" y2="{ty}" stroke="#666" stroke-width="2"/>')
                
                # 添加边标签
                if edge.get("label"):
                    label_x = (sx + tx) / 2
                    label_y = (sy + ty) / 2 - 10
                    svg_parts.append(f'<text x="{label_x}" y="{label_y}" text-anchor="middle" font-size="12" fill="#333">{edge["label"]}</text>')
        
        # 绘制节点
        for node in nodes:
            x = node.get("x", 0)
            y = node.get("y", 0)
            w = node.get("width", 120)
            h = node.get("height", 60)
            node_type = node.get("type", "process")
            label = node.get("label", "")
            color = node.get("color", self._get_node_color(node_type))
            
            # 根据节点类型绘制不同形状
            if node_type in ["start", "end"]:
                # 椭圆形
                rx = w / 2
                ry = h / 2
                svg_parts.append(f'<ellipse cx="{x + rx}" cy="{y + ry}" rx="{rx}" ry="{ry}" fill="{color}" stroke="#333" stroke-width="2"/>')
                svg_parts.append(f'<text x="{x + rx}" y="{y + ry + 4}" text-anchor="middle" font-size="14" fill="#fff">{label}</text>')
            
            elif node_type == "decision":
                # 菱形
                cx = x + w / 2
                cy = y + h / 2
                svg_parts.append(f'<polygon points="{cx},{y} {x + w},{cy} {cx},{y + h} {x},{cy}" fill="{color}" stroke="#333" stroke-width="2"/>')
                svg_parts.append(f'<text x="{cx}" y="{cy + 4}" text-anchor="middle" font-size="14" fill="#fff">{label}</text>')
            
            else:
                # 矩形
                svg_parts.append(f'<rect x="{x}" y="{y}" width="{w}" height="{h}" rx="5" fill="{color}" stroke="#333" stroke-width="2"/>')
                svg_parts.append(f'<text x="{x + w/2}" y="{y + h/2 + 4}" text-anchor="middle" font-size="14" fill="#fff">{label}</text>')
        
        svg_parts.append('</svg>')
        
        return "\n".join(svg_parts)
    
    def _get_node_color(self, node_type: str) -> str:
        """获取节点颜色"""
        color_map = {
            "start": "#27ae60",
            "end": "#e74c3c",
            "process": "#3498db",
            "decision": "#f39c12",
            "input_output": "#9b59b6",
            "connector": "#1abc9c",
            "subprocess": "#e67e22",
            "loop": "#95a5a6"
        }
        return color_map.get(node_type, "#3498db")
    
    async def _list_templates(self, **kwargs) -> Dict[str, Any]:
        """
        列出可用模板
        
        Returns:
            模板列表
        """
        templates = {
            "flowchart": {
                "types": [e.value for e in DiagramType],
                "layouts": [e.value for e in LayoutType],
                "themes": [e.value for e in ColorTheme],
                "element_types": [e.value for e in ElementType]
            },
            "floorplan": {
                "element_types": ["building", "road", "green_space", "water", "parking", "entrance", "boundary"]
            },
            "ppt": {
                "layouts": ["title_only", "title_content", "content_only", "comparison", "blank"],
                "themes": ["default", "blue", "green", "orange", "purple"],
                "transitions": ["fade", "slide", "zoom", "dissolve"]
            }
        }
        
        return {
            "success": True,
            "message": "获取模板列表成功",
            "data": templates
        }
    
    def _get_timestamp(self) -> str:
        """获取时间戳"""
        from datetime import datetime
        return datetime.now().isoformat()
    
    def get_agent_info(self) -> Dict[str, Any]:
        """获取智能体调用信息"""
        return {
            "name": self.name,
            "description": self.description,
            "category": self.category,
            "parameters": {
                "action": {
                    "type": "string",
                    "description": "操作类型",
                    "required": True,
                    "enum": ["create_flowchart", "create_floorplan", "generate_ppt", "render", "list_templates"]
                },
                "nodes": {
                    "type": "array",
                    "description": "流程图节点列表，每个节点包含 id, type, label, x, y, width, height, color",
                    "required": False
                },
                "edges": {
                    "type": "array",
                    "description": "流程图边列表，每个边包含 id, source, target, label",
                    "required": False
                },
                "elements": {
                    "type": "array",
                    "description": "平面图元素列表",
                    "required": False
                },
                "content": {
                    "type": "string",
                    "description": "PPT 内容描述（用于自动生成幻灯片）",
                    "required": False
                },
                "slides": {
                    "type": "array",
                    "description": "幻灯片列表",
                    "required": False
                },
                "output_path": {
                    "type": "string",
                    "description": "输出文件路径",
                    "required": False
                },
                "config": {
                    "type": "object",
                    "description": "配置参数",
                    "required": False
                },
                "diagram_data": {
                    "type": "object",
                    "description": "图表数据（render 操作时使用）",
                    "required": False
                },
                "format": {
                    "type": "string",
                    "description": "输出格式（svg, png, json）",
                    "required": False,
                    "default": "svg"
                }
            },
            "examples": [
                {
                    "input": {
                        "action": "create_flowchart",
                        "nodes": [
                            {"id": "n1", "type": "start", "label": "开始"},
                            {"id": "n2", "type": "process", "label": "处理任务"},
                            {"id": "n3", "type": "decision", "label": "是否完成?"},
                            {"id": "n4", "type": "end", "label": "结束"}
                        ],
                        "edges": [
                            {"id": "e1", "source": "n1", "target": "n2"},
                            {"id": "e2", "source": "n2", "target": "n3"},
                            {"id": "e3", "source": "n3", "target": "n4", "label": "是"},
                            {"id": "e4", "source": "n3", "target": "n2", "label": "否"}
                        ],
                        "config": {"title": "工作流程图", "layout": "top_down"}
                    },
                    "description": "创建一个简单的工作流程图"
                },
                {
                    "input": {
                        "action": "generate_ppt",
                        "content": "人工智能发展现状\n\n第一部分：AI 技术概述\n人工智能是研究、开发用于模拟、延伸和扩展人的智能的理论、方法、技术及应用系统的一门新技术科学。\n\n第二部分：应用场景\nAI 已广泛应用于医疗、金融、教育、交通等领域。\n\n第三部分：未来展望\n未来 AI 将在更多领域发挥重要作用。",
                        "output_path": "~/AI_presentation.pptx",
                        "config": {"title": "AI 技术介绍", "author": "演讲者"}
                    },
                    "description": "根据内容描述自动生成 PPT"
                },
                {
                    "input": {
                        "action": "list_templates"
                    },
                    "description": "列出所有可用模板"
                }
            ],
            "supported_diagram_types": [e.value for e in DiagramType],
            "supported_layouts": [e.value for e in LayoutType],
            "supported_themes": [e.value for e in ColorTheme],
            "supported_element_types": [e.value for e in ElementType]
        }


huashu_design_tool = HuashuDesignTool()


def get_huashu_design_tool() -> HuashuDesignTool:
    """获取华数设计工具实例"""
    return huashu_design_tool


async def test_huashu_design_tool():
    """测试华数设计工具"""
    import sys
    logger.remove()
    logger.add(sys.stdout, format="{time:YYYY-MM-DD HH:mm:ss} | {level} | {message}", colorize=False)
    
    print("=" * 60)
    print("测试 HuashuDesignTool")
    print("=" * 60)
    
    tool = HuashuDesignTool()
    
    # 测试列出模板
    print("\n[1] 测试列出模板:")
    result = await tool.execute(action="list_templates")
    if result["success"]:
        print("  ✓ 获取模板列表成功")
        templates = result["data"]
        print(f"  支持的流程图类型: {templates['flowchart']['types']}")
        print(f"  支持的布局: {templates['flowchart']['layouts']}")
        print(f"  支持的主题: {templates['flowchart']['themes']}")
    else:
        print(f"  ✗ 获取模板列表失败: {result['message']}")
    
    # 测试创建流程图
    print("\n[2] 测试创建流程图:")
    flowchart_result = await tool.execute(
        action="create_flowchart",
        nodes=[
            {"id": "n1", "type": "start", "label": "开始"},
            {"id": "n2", "type": "process", "label": "数据分析"},
            {"id": "n3", "type": "decision", "label": "数据有效?"},
            {"id": "n4", "type": "process", "label": "生成报告"},
            {"id": "n5", "type": "end", "label": "结束"}
        ],
        edges=[
            {"id": "e1", "source": "n1", "target": "n2"},
            {"id": "e2", "source": "n2", "target": "n3"},
            {"id": "e3", "source": "n3", "target": "n4", "label": "是"},
            {"id": "e4", "source": "n3", "target": "n2", "label": "否"},
            {"id": "e5", "source": "n4", "target": "n5"}
        ],
        config={"title": "数据分析流程图", "layout": "top_down", "theme": "blue"}
    )
    
    if flowchart_result["success"]:
        print("  ✓ 流程图创建成功")
        data = flowchart_result["data"]
        print(f"  节点数: {len(data['nodes'])}")
        print(f"  边数: {len(data['edges'])}")
        print(f"  图表ID: {flowchart_result['diagram_id']}")
        
        # 测试渲染图表
        print("\n[3] 测试渲染图表:")
        render_result = await tool.execute(
            action="render",
            diagram_data=data,
            format="svg"
        )
        
        if render_result["success"]:
            print(f"  ✓ 图表渲染成功")
            print(f"  输出路径: {render_result['output_path']}")
        else:
            print(f"  ✗ 图表渲染失败: {render_result['message']}")
    else:
        print(f"  ✗ 流程图创建失败: {flowchart_result['message']}")
    
    # 测试创建平面图
    print("\n[4] 测试创建平面图:")
    floorplan_result = await tool.execute(
        action="create_floorplan",
        elements=[
            {"id": "b1", "type": "building", "label": "主楼", "x": 200, "y": 150, "width": 300, "height": 200, "color": "#7f8c8d"},
            {"id": "r1", "type": "road", "label": "主干道", "x": 0, "y": 300, "width": 1200, "height": 30, "color": "#95a5a6"},
            {"id": "g1", "type": "green_space", "label": "绿地", "x": 550, "y": 150, "width": 200, "height": 150, "color": "#27ae60"}
        ],
        config={"title": "校园总平面图", "width": 1200, "height": 800, "scale": 1.0, "units": "米"}
    )
    
    if floorplan_result["success"]:
        print("  ✓ 平面图创建成功")
        print(f"  元素数: {len(floorplan_result['data']['elements'])}")
        print(f"  平面图ID: {floorplan_result['plan_id']}")
    else:
        print(f"  ✗ 平面图创建失败: {floorplan_result['message']}")
    
    # 测试生成 PPT
    print("\n[5] 测试生成 PPT:")
    ppt_result = await tool.execute(
        action="generate_ppt",
        content="华数设计工具介绍\n\n一、功能概述\n华数设计工具提供流程图绘制、总平面图绘制和 PPT 生成功能。\n\n二、技术特点\n支持多种图表类型，提供丰富的模板选择，支持自动布局。\n\n三、应用场景\n适用于企业汇报、项目管理、教育演示等多种场景。",
        config={"title": "华数设计工具演示", "author": "LivingTreeAI"}
    )
    
    if ppt_result["success"]:
        print("  ✓ PPT 生成成功")
        print(f"  幻灯片数: {ppt_result['slide_count']}")
        print(f"  输出路径: {ppt_result['output_path']}")
    else:
        print(f"  ✗ PPT 生成失败: {ppt_result['message']}")
    
    print("\n" + "=" * 60)
    print("测试完成")
    print("=" * 60)


if __name__ == "__main__":
    import asyncio
    asyncio.run(test_huashu_design_tool())