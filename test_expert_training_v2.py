"""
ExpertTrainingPipeline - 专家训练统一测试脚本 v2
================================================
优化版：修复了 DRET KB 共享实例 + 支持文件夹/压缩包/多格式文档训练

架构原则：
  1. 统一通过 Agent Chat 入口（HermesAgent.send_message）执行意图分析
  2. DRET 内嵌 KB 与外部 KB 共享同一实例（修复知识空白填充率0%问题）
  3. 支持文件夹递归扫描 + 压缩包（zip/7z/tar.gz等）自动解压
  4. 文档统一转换为 Markdown 再入知识库（docx/pdf/xlsx/html/txt 等）
  5. 支持知识库文档朗读（SoundEngine.play_voice_cn）

使用方式：
  python test_expert_training_v2.py                    # 完整流程
  python test_expert_training_v2.py --folder <path>     # 指定文件夹训练
  python test_expert_training_v2.py --archive <path>     # 压缩包训练
  python test_expert_training_v2.py --chat              # 仅 Agent Chat 测试
  python test_expert_training_v2.py --read-aloud <path> # 朗读指定文件
"""

import sys
import os
import time
import json
import tempfile
import shutil
import sqlite3
import uuid
import hashlib
from pathlib import Path
from typing import List, Dict, Any, Optional, Iterator
from dataclasses import dataclass, field
from concurrent.futures import ThreadPoolExecutor, as_completed

# ── 项目路径设置 ───────────────────────────────────────────────────────
PROJECT_ROOT = os.path.dirname(os.path.abspath(__file__))
if PROJECT_ROOT not in sys.path:
    sys.path.insert(0, PROJECT_ROOT)

# ════════════════════════════════════════════════════════════════════════
# 第一部分：文档格式转换器（多格式 → Markdown）
# ════════════════════════════════════════════════════════════════════════

class DocumentConverter:
    """
    统一文档格式转换器
    支持格式：txt, md, json, html, docx, xlsx, csv, pptx, pdf
    统一输出：Markdown 格式文本
    """

    # 支持的文件扩展名
    SUPPORTED_EXTENSIONS = {
        '.txt', '.md', '.markdown', '.json',
        '.html', '.htm', '.docx', '.xlsx', '.csv',
        '.pptx', '.pdf', '.xml', '.rtf'
    }

    # 文本类扩展名（直接读取）
    TEXT_EXTS = {'.txt', '.md', '.markdown', '.json', '.xml', '.rtf', '.html', '.htm', '.csv'}

    def __init__(self, max_file_size_mb: float = 50.0):
        self.max_file_size = max_file_size_mb * 1024 * 1024

    def is_supported(self, file_path: str) -> bool:
        """检查文件是否支持转换"""
        ext = Path(file_path).suffix.lower()
        return ext in self.SUPPORTED_EXTENSIONS

    def convert(self, file_path: str) -> Optional[Dict[str, Any]]:
        """
        转换文件为 Markdown 格式

        Returns:
            {"doc_id", "title", "content", "source_file", "format", "size"}
            or None if conversion failed
        """
        path = Path(file_path)

        if not path.exists():
            return None

        if path.stat().st_size > self.max_file_size:
            print(f"    [WARN] 文件过大，跳过: {path.name} ({path.stat().st_size / 1024 / 1024:.1f}MB)")
            return None

        ext = path.suffix.lower()
        title = path.stem

        # 文本类：直接读取
        if ext in self.TEXT_EXTS:
            return self._convert_text(path, title, ext)

        # Office/PDF 类
        if ext == '.docx':
            return self._convert_docx(path, title)
        elif ext == '.pdf':
            return self._convert_pdf(path, title)
        elif ext in {'.xlsx', '.xls'}:
            return self._convert_xlsx(path, title)
        elif ext == '.pptx':
            return self._convert_pptx(path, title)
        else:
            # 尝试当作文本读取
            return self._convert_text(path, title, ext)

    def _convert_text(self, path: Path, title: str, ext: str) -> Optional[Dict]:
        """转换文本类文件"""
        try:
            encodings = ['utf-8', 'gbk', 'gb2312', 'gb18030', 'latin-1']
            content = None

            for enc in encodings:
                try:
                    content = path.read_text(encoding=enc)
                    break
                except (UnicodeDecodeError, LookupError):
                    continue

            if content is None:
                return None

            # HTML 提取正文
            if ext in {'.html', '.htm'}:
                content = self._extract_html_text(content)

            return {
                "doc_id": self._gen_id(path),
                "title": title,
                "content": content.strip(),
                "source_file": str(path),
                "format": ext.lstrip('.'),
                "size": path.stat().st_size
            }
        except Exception as e:
            print(f"    [WARN] 文本转换失败 [{path.name}]: {e}")
            return None

    def _extract_html_text(self, html: str) -> str:
        """从 HTML 提取纯文本"""
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(html, 'html.parser')
            # 移除 script 和 style
            for tag in soup(['script', 'style', 'noscript']):
                tag.decompose()
            text = soup.get_text(separator='\n', strip=True)
            # 清理多余空行
            import re
            text = re.sub(r'\n{3,}', '\n\n', text)
            return text
        except ImportError:
            # fallback: 正则提取
            import re
            text = re.sub(r'<[^>]+>', ' ', html)
            text = re.sub(r'\s{2,}', ' ', text)
            return text.strip()

    def _convert_docx(self, path: Path, title: str) -> Optional[Dict]:
        """转换 Word 文档"""
        try:
            import docx
            doc = docx.Document(str(path))
            paragraphs = []

            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    # 标题检测（粗体或特定样式）
                    if para.style and 'Heading' in str(para.style.name):
                        paragraphs.append(f"## {text}")
                    elif para.text.startswith('#'):
                        paragraphs.append(text)
                    else:
                        paragraphs.append(text)

            # 提取表格
            for table in doc.tables:
                rows_data = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows_data.append(' | '.join(cells))
                if rows_data:
                    paragraphs.append('\n'.join(['| ' + r + ' |' for r in rows_data[:5]]))

            content = '\n\n'.join(paragraphs)
            return {
                "doc_id": self._gen_id(path),
                "title": title,
                "content": content,
                "source_file": str(path),
                "format": "docx",
                "size": path.stat().st_size
            }
        except ImportError:
            print("    [WARN] python-docx 未安装，尝试 docx2txt...")
            return self._convert_docx_fallback(path, title)
        except Exception as e:
            print(f"    [WARN] DOCX 转换失败 [{path.name}]: {e}")
            return None

    def _convert_docx_fallback(self, path: Path, title: str) -> Optional[Dict]:
        """docx2txt 后备方案"""
        try:
            import docx2txt
            content = docx2txt.process(str(path))
            return {
                "doc_id": self._gen_id(path),
                "title": title,
                "content": content.strip(),
                "source_file": str(path),
                "format": "docx",
                "size": path.stat().st_size
            }
        except ImportError:
            return None
        except Exception as e:
            return None

    def _convert_pdf(self, path: Path, title: str) -> Optional[Dict]:
        """转换 PDF 文档"""
        extractors = [
            ('pymupdf', self._extract_pdf_pymupdf),
            ('pdfplumber', self._extract_pdf_plumber),
            ('pypdf2', self._extract_pdf_pypdf2),
        ]

        for name, extractor in extractors:
            try:
                result = extractor(path, title)
                if result:
                    return result
            except ImportError:
                continue
            except Exception as e:
                print(f"    [WARN] {name} 提取失败 [{path.name}]: {e}")
                continue

        print(f"    [WARN] PDF 缺少依赖（PyMuPDF/pdfplumber/PyPDF2），跳过: {path.name}")
        return None

    def _extract_pdf_pymupdf(self, path: Path, title: str) -> Optional[Dict]:
        import fitz
        doc = fitz.open(str(path))
        paragraphs = []
        for page_num, page in enumerate(doc, 1):
            text = page.get_text()
            if text.strip():
                paragraphs.append(f"### 第 {page_num} 页\n{text.strip()}")
        doc.close()
        content = '\n\n'.join(paragraphs)
        return {
            "doc_id": self._gen_id(path),
            "title": title,
            "content": content,
            "source_file": str(path),
            "format": "pdf",
            "size": path.stat().st_size
        }

    def _extract_pdf_plumber(self, path: Path, title: str) -> Optional[Dict]:
        import pdfplumber
        paragraphs = []
        with pdfplumber.open(str(path)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ''
                if text.strip():
                    paragraphs.append(f"### 第 {i} 页\n{text.strip()}")
        content = '\n\n'.join(paragraphs)
        return {
            "doc_id": self._gen_id(path),
            "title": title,
            "content": content,
            "source_file": str(path),
            "format": "pdf",
            "size": path.stat().st_size
        }

    def _extract_pdf_pypdf2(self, path: Path, title: str) -> Optional[Dict]:
        import PyPDF2
        paragraphs = []
        with open(str(path), 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ''
                if text.strip():
                    paragraphs.append(f"### 第 {i} 页\n{text.strip()}")
        content = '\n\n'.join(paragraphs)
        return {
            "doc_id": self._gen_id(path),
            "title": title,
            "content": content,
            "source_file": str(path),
            "format": "pdf",
            "size": path.stat().st_size
        }

    def _convert_xlsx(self, path: Path, title: str) -> Optional[Dict]:
        """转换 Excel 文档"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(path), data_only=True)
            paragraphs = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                paragraphs.append(f"### 工作表: {sheet_name}")
                rows = list(ws.iter_rows(values_only=True))
                # 最多取前100行
                for row in rows[:100]:
                    cells = [str(c) if c is not None else '' for c in row]
                    line = ' | '.join(cells).strip()
                    if line and not line.startswith('|'):
                        paragraphs.append('| ' + ' | '.join(cells) + ' |')

                if len(wb.sheetnames) > 3:
                    paragraphs.append(f"... 共 {len(wb.sheetnames)} 个工作表")
                    break

            content = '\n'.join(paragraphs)
            return {
                "doc_id": self._gen_id(path),
                "title": title,
                "content": content,
                "source_file": str(path),
                "format": "xlsx",
                "size": path.stat().st_size
            }
        except ImportError:
            return self._convert_csv_fallback(path, title)
        except Exception as e:
            print(f"    [WARN] XLSX 转换失败 [{path.name}]: {e}")
            return None

    def _convert_pptx(self, path: Path, title: str) -> Optional[Dict]:
        """转换 PowerPoint 文档"""
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            paragraphs = []

            for i, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    paragraphs.append(f"### 幻灯片 {i}\n" + '\n'.join(slide_texts))

            content = '\n\n'.join(paragraphs)
            return {
                "doc_id": self._gen_id(path),
                "title": title,
                "content": content,
                "source_file": str(path),
                "format": "pptx",
                "size": path.stat().st_size
            }
        except ImportError:
            return None
        except Exception as e:
            print(f"    [WARN] PPTX 转换失败 [{path.name}]: {e}")
            return None

    def _convert_csv_fallback(self, path: Path, title: str) -> Optional[Dict]:
        """CSV 后备方案"""
        try:
            import csv
            rows = []
            with open(str(path), newline='', encoding='utf-8') as f:
                reader = csv.reader(f)
                for i, row in enumerate(reader):
                    if i >= 100:
                        break
                    rows.append('| ' + ' | '.join(row) + ' |')

            content = '\n'.join(rows)
            return {
                "doc_id": self._gen_id(path),
                "title": title,
                "content": content,
                "source_file": str(path),
                "format": "csv",
                "size": path.stat().st_size
            }
        except Exception:
            return None

    @staticmethod
    def _gen_id(path: Path) -> str:
        """生成文档 ID"""
        return hashlib.md5(str(path).encode()).hexdigest()[:12]


# ════════════════════════════════════════════════════════════════════════
# 第二部分：归档文件处理器（支持 zip/7z/tar.gz 等）
# ════════════════════════════════════════════════════════════════════════

class ArchiveProcessor:
    """
    归档文件自动处理器
    支持格式：.zip, .7z, .tar, .tar.gz, .tar.bz2, .tar.xz, .rar
    """

    # 压缩包扩展名
    ARCHIVE_EXTS = {
        '.zip', '.7z', '.tar', '.gz', '.bz2', '.xz',
        '.tar.gz', '.tgz', '.tar.bz2', '.tar.xz', '.rar'
    }

    def __init__(self, work_dir: Optional[str] = None):
        self.work_dir = work_dir or tempfile.mkdtemp(prefix="hermes_archive_")

    def is_archive(self, file_path: str) -> bool:
        ext = Path(file_path).suffix.lower()
        if ext in {'.gz', '.bz2', '.xz'} and not file_path.endswith('.tar.gz'):
            return True
        return ext in self.ARCHIVE_EXTS

    def extract(self, archive_path: str) -> str:
        """
        解压归档文件到临时目录

        Returns:
            解压后的目录路径
        """
        path = Path(archive_path)
        extract_dir = os.path.join(self.work_dir, path.stem.replace('.', '_'))
        os.makedirs(extract_dir, exist_ok=True)

        print(f"  [i] 解压归档: {path.name} → {extract_dir}")

        if self._try_nanazip(archive_path, extract_dir):
            return extract_dir

        if self._try_7z(archive_path, extract_dir):
            return extract_dir

        if self._try_python(archive_path, extract_dir):
            return extract_dir

        print(f"  [WARN] 所有解压方式均失败: {archive_path}")
        return extract_dir

    def _try_nanazip(self, archive_path: str, out_dir: str) -> bool:
        """使用 NanaZip (7zip fork) 解压"""
        nanazip = shutil.which("7z") or shutil.which("7za")
        if not nanazip:
            return False
        try:
            result = shutil.run(
                [nanazip, "x", "-y", "-o" + out_dir, archive_path],
                capture_output=True, timeout=300
            )
            return result.returncode == 0
        except Exception:
            return False

    def _try_7z(self, archive_path: str, out_dir: str) -> bool:
        """使用 Python zipfile / tarfile"""
        path = Path(archive_path)
        ext = path.suffix.lower()

        try:
            if ext == '.zip':
                shutil.unpack_archive(archive_path, out_dir, 'zip')
                return True
            elif ext in {'.tar', '.tar.gz', '.tgz', '.tar.bz2', '.tar.xz'}:
                shutil.unpack_archive(archive_path, out_dir, 'tar')
                return True
            elif ext in {'.gz', '.bz2', '.xz'}:
                # 单文件压缩包
                base = Path(out_dir) / path.stem
                with open(base, 'wb') as out:
                    with open(archive_path, 'rb') as inp:
                        out.write(self._decompress(archive_path))
                return True
        except Exception as e:
            print(f"  [WARN] Python 解压失败: {e}")
        return False

    def _decompress(self, path: str) -> bytes:
        import gzip, bz2, lzma
        p = Path(path)
        with open(path, 'rb') as f:
            data = f.read()
        if p.suffix == '.gz':
            return gzip.decompress(data)
        elif p.suffix == '.bz2':
            return bz2.decompress(data)
        elif p.suffix == '.xz':
            return lzma.decompress(data)
        return data

    def _try_python(self, archive_path: str, out_dir: str) -> bool:
        """Python shutil 后备"""
        try:
            shutil.unpack_archive(archive_path, out_dir)
            return True
        except Exception:
            return False

    def scan_directory(self, directory: str, converter: DocumentConverter,
                       max_depth: int = 5) -> List[Dict[str, Any]]:
        """
        递归扫描目录，收集所有可转换文档

        Args:
            directory: 扫描目录
            converter: 文档转换器
            max_depth: 最大递归深度

        Returns:
            文档列表 [{"doc_id", "title", "content", ...}]
        """
        docs = []
        scanned = 0
        skipped_archive = 0
        failed = 0

        def _scan(dir_path: str, depth: int):
            nonlocal scanned, skipped_archive, failed

            if depth > max_depth:
                return

            try:
                entries = list(Path(dir_path).iterdir())
            except PermissionError:
                return

            for entry in entries:
                if entry.is_file():
                    scanned += 1
                    str_path = str(entry)

                    # 跳过临时文件
                    if entry.name.startswith('.') or '~$' in entry.name:
                        continue

                    # 压缩包：递归解压扫描
                    if ArchiveProcessor.is_archive(str_path):
                        try:
                            sub_dir = self.extract(str_path)
                            _scan(sub_dir, depth + 1)
                            skipped_archive += 1
                        except Exception:
                            pass
                        continue

                    # 支持的文档格式
                    if converter.is_supported(str_path):
                        doc = converter.convert(str_path)
                        if doc and len(doc['content'].strip()) > 50:
                            docs.append(doc)
                            print(f"    [OK] {entry.name} ({doc['format']}, {len(doc['content'])}字)")
                        else:
                            failed += 1

                elif entry.is_dir():
                    _scan(str(entry), depth + 1)

        print(f"  [i] 开始扫描目录: {directory}")
        _scan(directory, 0)
        print(f"  [i] 扫描完成: {scanned} 文件, {len(docs)} 文档, {skipped_archive} 归档, {failed} 失败")

        return docs

    def cleanup(self):
        """清理临时解压目录"""
        if os.path.exists(self.work_dir):
            shutil.rmtree(self.work_dir, ignore_errors=True)


# ════════════════════════════════════════════════════════════════════════
# 第三部分：知识库文档朗读器（SoundEngine）
# ════════════════════════════════════════════════════════════════════════

class KnowledgeBaseReader:
    """
    知识库文档朗读器
    使用 Windows SAPI 中文语音或 edge-tts 进行朗读
    """

    def __init__(self):
        self.speaker = None
        self._init_speaker()

    def _init_speaker(self):
        """初始化 Windows SAPI 语音"""
        try:
            import win32com.client
            import pythoncom
            pythoncom.CoInitialize()
            self.speaker = win32com.client.Dispatch("SAPI.SpVoice")
            print("[TTS] Windows SAPI 语音引擎初始化成功")
        except ImportError:
            print("[TTS] pywin32 未安装，将使用 edge-tts 后备方案")
            self.speaker = None
        except Exception as e:
            print(f"[TTS] SAPI 初始化失败: {e}")
            self.speaker = None

    def speak(self, text: str, rate: int = 0, use_cn_voice: bool = True) -> bool:
        """
        朗读文本

        Args:
            text: 要朗读的文本
            rate: 语速 (-10 到 10)
            use_cn_voice: 是否使用中文语音
        """
        if not text or not self.speaker:
            return False

        try:
            if use_cn_voice:
                return self._speak_cn(text, rate)
            else:
                self.speaker.Rate = max(-10, min(10, rate))
                self.speaker.Speak(text)
                return True
        except Exception as e:
            print(f"[TTS] 朗读失败: {e}")
            return False

    def _speak_cn(self, text: str, rate: int = 0) -> bool:
        """使用中文语音朗读"""
        try:
            speaker = win32com.client.Dispatch("SAPI.SpVoice")
            voices = speaker.GetVoices()

            cn_voice = None
            for voice in voices:
                desc = voice.GetDescription()
                if 'Chinese' in desc or '中文' in desc:
                    cn_voice = voice
                    break

            if cn_voice:
                speaker.Voice = cn_voice

            if rate != 0:
                speaker.Rate = max(-10, min(10, rate))

            speaker.Speak(text)
            return True
        except Exception as e:
            print(f"[TTS] 中文语音朗读失败: {e}")
            return False

    async def speak_async(self, text: str) -> bool:
        """异步朗读（使用 edge-tts）"""
        try:
            import edge_tts
            from edge_tts import Communicate

            audio_file = tempfile.NamedTemporaryFile(suffix='.mp3', delete=False)
            audio_path = audio_file.name
            audio_file.close()

            comm = Communicate(text, voice="zh-CN-XiaoxiaoNeural")
            await comm.save(audio_path)

            # 播放音频
            if os.name == 'nt':
                import winsound
                winsound.PlaySound(audio_path, winsound.SND_FILENAME)

            os.unlink(audio_path)
            return True
        except ImportError:
            print("[TTS] edge-tts 未安装，使用同步朗读")
            return self.speak(text)
        except Exception as e:
            print(f"[TTS] edge-tts 异步朗读失败: {e}")
            return False

    def read_document(self, doc: Dict[str, Any], max_chars: int = 2000) -> bool:
        """
        朗读知识库文档

        Args:
            doc: 知识库文档
            max_chars: 最大朗读字符数
        """
        title = doc.get('title', '未知文档')
        content = doc.get('content', '')[:max_chars]

        print(f"\n  [TTS] 开始朗读: {title}")
        print(f"  ────────────")
        print(f"  {content[:300]}...")
        print(f"  ────────────")

        return self.speak(f"{title}。{content}", rate=0)


# ════════════════════════════════════════════════════════════════════════
# 第四部分：专家训练统一管道（Agent Chat 单一入口）
# ════════════════════════════════════════════════════════════════════════

class ExpertTrainingPipeline:
    """
    专家训练统一管道 - Agent Chat 单一入口

    整合所有组件：
    - HermesAgent（统一聊天入口）
    - KnowledgeBaseLayer（知识库）
    - DRET 系统（专家训练）
    - ArchiveProcessor（文件夹/压缩包处理）
    - DocumentConverter（多格式文档转换）
    - KnowledgeBaseReader（朗读功能）

    架构原则：
    1. 所有用户交互通过 HermesAgent.send_message() 统一入口
    2. DRET KB 与外部 KB 共享同一实例（修复填充率0%问题）
    3. 支持文件夹/压缩包自动训练
    4. 文档统一转 MD 入库
    5. 支持知识库文档朗读
    """

    def __init__(self, session_id: str = "expert_training_session"):
        self.session_id = session_id
        self.agent = None
        self.kb = None
        self.dret = None
        self.archive_processor = None
        self.converter = DocumentConverter()
        self.reader = KnowledgeBaseReader()
        self._initialized = False
        self._init_components()

    def _init_components(self):
        """初始化所有组件"""
        print("=" * 70)
        print("  ExpertTrainingPipeline 初始化")
        print("=" * 70)

        # 1. 知识库（最优先，确保 DRET 共享）
        print("\n  [1/5] 初始化知识库...")
        try:
            from core.fusion_rag.knowledge_base import KnowledgeBaseLayer
            self.kb = KnowledgeBaseLayer()
            print(f"  [OK] KnowledgeBaseLayer 嵌入模型: BAAI/bge-small-zh")
        except Exception as e:
            print(f"  [WARN] KnowledgeBaseLayer 失败: {e}")
            self.kb = None

        # 2. DRET 系统（★ 关键：与 KB 共享实例）
        print("\n  [2/5] 初始化 DRET 专家训练系统...")
        try:
            from core.skill_evolution.dret_l04_integration import create_l04_dret_system

            self.dret = create_l04_dret_system(
                max_recursion_depth=5,
                enable_l04=True,
                enable_expert=True
            )

            # ★★★ 修复 KB 共享实例 ★★★
            if self.kb is not None and hasattr(self.dret, 'gap_detector'):
                old_kb = getattr(self.dret.gap_detector, 'knowledge_base', None)
                self.dret.gap_detector.knowledge_base = self.kb
                print(f"  [OK] DRET KB 共享实例已修复（旧实例: {old_kb is not None} → 新实例: {self.kb}）")
            else:
                print(f"  [WARN] DRET gap_detector 不存在，跳过 KB 共享")

            print(f"  [OK] DRET 递归深度: {self.dret.max_depth}")
        except Exception as e:
            print(f"  [WARN] DRET 初始化失败: {e}")
            import traceback
            traceback.print_exc()
            self.dret = None

        # 3. HermesAgent（Agent Chat 统一入口）
        print("\n  [3/5] 初始化 HermesAgent（Agent Chat 入口）...")
        try:
            from core.config import AppConfig
            from client.src.business.agent import HermesAgent, AgentCallbacks

            config = AppConfig()
            callbacks = AgentCallbacks(
                stream_delta=lambda delta: print(delta, end='', flush=True),
                stats_update=lambda stats: None
            )
            self.agent = HermesAgent(config, session_id=self.session_id, callbacks=callbacks)
            print(f"  [OK] HermesAgent 已初始化，会话ID: {self.session_id}")
        except Exception as e:
            print(f"  [WARN] HermesAgent 初始化失败: {e}")
            import traceback
            traceback.print_exc()
            self.agent = None

        # 4. 归档处理器
        print("\n  [4/5] 初始化归档处理器...")
        self.archive_processor = ArchiveProcessor()
        print(f"  [OK] ArchiveProcessor 工作目录: {self.archive_processor.work_dir}")

        # 5. 朗读器
        print("\n  [5/5] 初始化朗读器...")
        print(f"  [OK] KnowledgeBaseReader 初始化完成")

        self._initialized = True
        print("\n" + "=" * 70)
        print("  ✅ 初始化完成")
        print("=" * 70)

    # ── Agent Chat 统一入口 ──────────────────────────────────────────────

    def chat(self, message: str) -> str:
        """
        通过 HermesAgent 发送消息（统一入口）

        Args:
            message: 用户消息

        Returns:
            Agent 响应文本
        """
        if not self.agent:
            return "[ERROR] HermesAgent 未初始化"

        print(f"\n{'─' * 70}")
        print(f"  [Agent Chat] {message[:80]}{'...' if len(message) > 80 else ''}")
        print(f"{'─' * 70}")

        try:
            response_chunks = []
            for chunk in self.agent.send_message(message):
                if hasattr(chunk, 'delta') and chunk.delta:
                    response_chunks.append(chunk.delta)
                elif hasattr(chunk, 'error') and chunk.error:
                    return f"[ERROR] {chunk.error}"

            response = ''.join(response_chunks)
            print(f"\n  [Agent Chat] 响应长度: {len(response)} 字符")
            return response

        except sqlite3.IntegrityError as e:
            # FOREIGN KEY 错误：会话不存在，尝试重建会话
            print(f"  [WARN] Session不存在，重建会话...")
            new_id = f"exp_{uuid.uuid4().hex[:8]}"
            try:
                self.agent.session_db.create_session(model="qwen2.5:1.5b", session_id=new_id)
                self.agent.session_id = new_id
                return self.chat(message)
            except Exception:
                return f"[ERROR] 会话重建失败: {e}"
        except sqlite3.OperationalError as e:
            # 事务冲突：创建新会话
            if "transaction" in str(e):
                print(f"  [WARN] 事务冲突，创建新会话...")
                new_id = f"exp_{uuid.uuid4().hex[:8]}"
                try:
                    self.agent.session_db.create_session(model="qwen2.5:1.5b", session_id=new_id)
                    self.agent.session_id = new_id
                    return self.chat(message)
                except Exception:
                    return f"[ERROR] 会话创建失败: {e}"
            return f"[ERROR] {e}"
        except Exception as e:
            import traceback
            traceback.print_exc()
            return f"[ERROR] {e}"

    # ── 文档入知识库 ──────────────────────────────────────────────────────

    def ingest_documents(self, documents: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        向知识库批量注入文档（同时注入 pipeline.kb 和 HermesAgent.kb）

        Args:
            documents: 文档列表 [{"id", "title", "content", ...}]

        Returns:
            注入结果统计
        """
        stats = {"total": len(documents), "success": 0, "failed": 0, "total_chars": 0}

        for doc in documents:
            try:
                # 文档必须有 id, title, content
                if 'id' not in doc:
                    doc['id'] = hashlib.md5(doc.get('title', '').encode()).hexdigest()[:12]

                # 注入到 pipeline.kb（DRET 共享实例）
                if self.kb:
                    self.kb.add_document(doc)

                # 注入到 HermesAgent.kb（KnowledgeBaseVectorStore）
                if self.agent and hasattr(self.agent, 'knowledge_base'):
                    try:
                        self.agent.knowledge_base.add_knowledge(
                            doc_id=doc['id'],
                            content=doc['content'],
                            metadata={"title": doc.get('title', ''), "source": "expert_training"}
                        )
                    except AttributeError:
                        # KnowledgeBaseVectorStore 没有 add_knowledge，跳过
                        pass

                stats["success"] += 1
                stats["total_chars"] += len(doc.get('content', ''))

            except Exception as e:
                stats["failed"] += 1
                print(f"    [WARN] 文档注入失败 [{doc.get('title', '?')}]: {e}")

        print(f"\n  [i] 文档注入完成: {stats['success']}/{stats['total']} 成功, "
              f"{stats['total_chars']:,} 字符入库")
        return stats

    def ingest_folder(self, folder_path: str, max_depth: int = 5) -> Dict[str, Any]:
        """
        从文件夹批量导入文档到知识库

        Args:
            folder_path: 文件夹路径
            max_depth: 最大递归深度

        Returns:
            导入统计
        """
        print(f"\n{'─' * 70}")
        print(f"  [知识库导入] 文件夹: {folder_path}")
        print(f"{'─' * 70}")

        docs = self.archive_processor.scan_directory(
            folder_path, self.converter, max_depth
        )

        if not docs:
            print(f"  [WARN] 未找到可导入的文档")
            return {"success": False, "found": 0}

        return self.ingest_documents(docs)

    def ingest_archive(self, archive_path: str) -> Dict[str, Any]:
        """
        从压缩包导入文档到知识库

        Args:
            archive_path: 压缩包路径

        Returns:
            导入统计
        """
        print(f"\n{'─' * 70}")
        print(f"  [知识库导入] 压缩包: {archive_path}")
        print(f"{'─' * 70}")

        extract_dir = self.archive_processor.extract(archive_path)
        return self.ingest_folder(extract_dir, max_depth=5)

    # ── 专家训练 ──────────────────────────────────────────────────────────

    def expert_training(self, doc_content: str, doc_id: str,
                        recursion_depth: int = 5) -> Optional[Dict]:
        """
        执行专家训练（通过 DRET 系统）

        Args:
            doc_content: 训练文档内容
            doc_id: 文档ID
            recursion_depth: 递归深度

        Returns:
            训练报告
        """
        if not self.dret:
            return None

        print(f"\n{'─' * 70}")
        print(f"  [专家训练] 文档: {doc_id}, 递归深度: {recursion_depth}")
        print(f"{'─' * 70}")

        t0 = time.time()
        try:
            report = self.dret.learn_from_document(
                doc_content=doc_content,
                doc_id=doc_id,
                session_id=self.session_id,
                recursion_depth=recursion_depth
            )
            elapsed = time.time() - t0

            # 打印训练摘要
            print(f"\n  ┌─ DRET 训练报告")
            print(f"  │  文档ID      : {report.get('doc_id')}")
            print(f"  │  递归深度    : {report.get('max_depth_used')} 层")
            print(f"  │  知识空白发现: {report.get('gaps_found')} 个")
            print(f"  │  知识空白填充: {report.get('gaps_filled')} 个")
            print(f"  │  矛盾发现    : {report.get('conflicts_found')} 个")
            print(f"  │  知识图谱    : {report.get('knowledge_graph', {}).get('nodes', 0)} 节点, "
                  f"{report.get('knowledge_graph', {}).get('edges', 0)} 边")
            print(f"  │  总耗时      : {elapsed:.2f}s")
            if report.get('expert_persona'):
                print(f"  │  专家人格    : {report['expert_persona']}")
            print(f"  └─")

            # ★ KB 共享修复效果验证
            fill_rate = report.get('gaps_filled', 0) / max(report.get('gaps_found', 1), 1) * 100
            if fill_rate > 0:
                print(f"\n  ✅ KB共享修复成功！知识填充率: {fill_rate:.1f}%")
            else:
                print(f"\n  ⚠ KB共享修复效果待验证，填充率: {fill_rate:.1f}%")

            return report

        except Exception as e:
            import traceback
            traceback.print_exc()
            print(f"  [FAIL] 专家训练失败: {e}")
            return None

    def expert_training_folder(self, folder_path: str, recursion_depth: int = 5) -> List[Dict]:
        """
        对文件夹中所有文档执行专家训练

        Args:
            folder_path: 文件夹路径
            recursion_depth: 递归深度

        Returns:
            训练报告列表
        """
        print(f"\n{'=' * 70}")
        print(f"  专家训练：批量文件夹模式")
        print(f"  目录: {folder_path}, 递归深度: {recursion_depth}")
        print(f"{'=' * 70}")

        # 1. 收集文档
        docs = self.archive_processor.scan_directory(
            folder_path, self.converter, max_depth=5
        )

        if not docs:
            print(f"  [WARN] 未找到可训练文档")
            return []

        # 2. 批量注入知识库
        self.ingest_documents(docs)

        # 3. 逐篇专家训练
        reports = []
        for doc in docs:
            report = self.expert_training(
                doc_content=doc['content'],
                doc_id=doc['id'],
                recursion_depth=recursion_depth
            )
            if report:
                reports.append(report)

        return reports

    # ── 知识库搜索 ────────────────────────────────────────────────────────

    def search_knowledge(self, query: str, top_k: int = 5) -> List[Dict]:
        """
        搜索知识库

        Args:
            query: 搜索关键词
            top_k: 返回结果数

        Returns:
            搜索结果列表
        """
        if not self.kb:
            return []

        results = self.kb.search(query, top_k=top_k)
        print(f"\n{'─' * 70}")
        print(f"  [知识库搜索] \"{query}\" → {len(results)} 条结果")
        print(f"{'─' * 70}")

        for i, r in enumerate(results, 1):
            print(f"  #{i} [{r.get('doc_id', '?')}] {r.get('title', '')[:30]}")
            print(f"      分数: {r.get('score', 0):.4f} | "
                  f"{r.get('content', '')[:100]}...")

        return results

    # ── 朗读 ──────────────────────────────────────────────────────────────

    def read_aloud_file(self, file_path: str, max_chars: int = 3000) -> bool:
        """
        朗读指定文件

        Args:
            file_path: 文件路径
            max_chars: 最大朗读字符数
        """
        print(f"\n{'─' * 70}")
        print(f"  [朗读] 文件: {file_path}")
        print(f"{'─' * 70}")

        if not os.path.exists(file_path):
            print(f"  [ERROR] 文件不存在: {file_path}")
            return False

        # 转换文件
        doc = self.converter.convert(file_path)
        if not doc:
            print(f"  [ERROR] 文件转换失败: {file_path}")
            return False

        print(f"  [OK] 文件转换成功: {doc['title']} ({doc['format']}, {len(doc['content'])}字)")

        # 朗读
        return self.reader.read_document(doc, max_chars=max_chars)

    def read_aloud_kb_doc(self, doc_id: str, max_chars: int = 3000) -> bool:
        """
        朗读知识库中的文档

        Args:
            doc_id: 文档ID
            max_chars: 最大朗读字符数
        """
        if not self.kb:
            print("[ERROR] 知识库未初始化")
            return False

        # 搜索该文档
        results = self.kb.search(doc_id, top_k=1)
        if not results:
            print(f"[ERROR] 知识库中未找到文档: {doc_id}")
            return False

        return self.reader.read_document(results[0], max_chars=max_chars)

    # ── 清理 ──────────────────────────────────────────────────────────────

    def cleanup(self):
        """清理资源"""
        if self.archive_processor:
            self.archive_processor.cleanup()
        if self.agent:
            self.agent.close()


# ════════════════════════════════════════════════════════════════════════
# 第五部分：CLI 入口 & 主流程
# ════════════════════════════════════════════════════════════════════════

def print_banner():
    print("""
╔══════════════════════════════════════════════════════════════════╗
║                                                                  ║
║   ExpertTrainingPipeline v2 - 专家训练系统                       ║
║   Agent Chat 统一入口 | DRET KB共享 | 多格式文档 | 朗读功能       ║
║                                                                  ║
╚══════════════════════════════════════════════════════════════════╝
    """)


def main():
    import argparse

    parser = argparse.ArgumentParser(description="ExpertTrainingPipeline v2")
    parser.add_argument('--folder', type=str, help='指定文件夹进行专家训练')
    parser.add_argument('--archive', type=str, help='指定压缩包进行专家训练')
    parser.add_argument('--chat', action='store_true', help='仅测试 Agent Chat')
    parser.add_argument('--read-aloud', type=str, dest='read_aloud', help='朗读指定文件')
    parser.add_argument('--depth', type=int, default=5, help='递归深度（默认5层）')
    parser.add_argument('--ingest', type=str, help='仅导入文档到知识库（文件夹路径）')
    args = parser.parse_args()

    print_banner()

    # 初始化管道
    pipeline = ExpertTrainingPipeline()

    # ── 模式1：仅 Agent Chat 测试 ─────────────────────────────────────
    if args.chat:
        print("\n" + "=" * 70)
        print("  模式：Agent Chat 统一入口测试")
        print("=" * 70)

        # 通过 Agent Chat 执行意图分析
        test_queries = [
            "帮我分析一下环评报告的结构",
            "搜索知识库中的污染物相关内容",
        ]

        for query in test_queries:
            response = pipeline.chat(query)
            print(f"\n  [响应]\n  {response[:500]}{'...' if len(response) > 500 else ''}\n")

    # ── 模式2：仅朗读测试 ─────────────────────────────────────────────
    elif args.read_aloud:
        print("\n" + "=" * 70)
        print("  模式：朗读测试")
        print("=" * 70)
        pipeline.read_aloud_file(args.read_aloud, max_chars=3000)

    # ── 模式3：仅导入文档 ─────────────────────────────────────────────
    elif args.ingest:
        print("\n" + "=" * 70)
        print("  模式：文档导入知识库")
        print("=" * 70)
        result = pipeline.ingest_folder(args.ingest, max_depth=5)
        print(f"\n  [结果] {result}")

    # ── 模式4：压缩包专家训练 ─────────────────────────────────────────
    elif args.archive:
        print("\n" + "=" * 70)
        print("  模式：压缩包专家训练")
        print("=" * 70)

        # 导入文档
        result = pipeline.ingest_archive(args.archive)
        print(f"  [导入结果] {result}")

        # 专家训练
        if result.get('success') and result.get('found', 0) > 0:
            reports = pipeline.expert_training_folder(args.archive, args.depth)
            print(f"\n  [训练完成] 共 {len(reports)} 份训练报告")

        # 知识库搜索测试
        pipeline.search_knowledge("污染物")

    # ── 模式5：文件夹专家训练（默认） ────────────────────────────────
    elif args.folder:
        print("\n" + "=" * 70)
        print("  模式：文件夹专家训练")
        print("=" * 70)

        reports = pipeline.expert_training_folder(args.folder, args.depth)
        print(f"\n  [训练完成] 共 {len(reports)} 份训练报告")

        # 知识库搜索测试
        pipeline.search_knowledge("污染物")

    # ── 默认：完整演示流程 ─────────────────────────────────────────────
    else:
        print("\n" + "=" * 70)
        print("  模式：完整演示流程（5步）")
        print("=" * 70)

        # Step 1: 预注入环评文档（确保 KB 有内容供 DRET 填充）
        print("\n\n>>> Step 0: 预注入环评文档（为专家训练提供背景知识）")
        print("─" * 70)
        EIA_PRE_DOCS = [
            {
                "id": "eia_pre_001",
                "title": "污染物排放标准",
                "content": "大气污染物排放必须满足《大气污染物综合排放标准》（GB16297）。"
                          "主要污染物包括SO2、NOx、PM2.5、VOCs等。废水排放须达到《污水综合排放标准》（GB8978）。"
                          "COD、氨氮、SO2、NOx为主要总量控制指标。"
            },
            {
                "id": "eia_pre_002",
                "title": "环境风险评价方法",
                "content": "采用高斯烟团模型预测有毒气体泄漏扩散，确定毒害区（IDLH范围）。"
                          "采用Streeter-Phelps模型进行水体预测。应急预案包含预警、疏散、消防、救援程序。"
            }
        ]
        pipeline.ingest_documents(EIA_PRE_DOCS)

        # Step 1: Agent Chat 意图分析
        print("\n\n>>> Step 1: Agent Chat 意图分析")
        print("─" * 70)
        response = pipeline.chat("我需要分析一份环评报告，它包含污染物排放和环境风险评价内容")
        print(f"\n意图分析响应: {response[:300]}...")

        # Step 2: 内置环评文档训练
        print("\n\n>>> Step 2: 环评报告专家训练（递归5层）")
        print("─" * 70)

        ENV_TRAINING_DOC = """
        环评报告编写要点：

        首先确定评价等级，然后开展现状调查。工程分析是报告核心，
        需要核算所有污染物的产生量和排放量。大气污染物包括SO2、NOx、颗粒物等。
        水污染物包括COD、氨氮、重金属等。

        基于高斯扩散模型进行大气预测。基于Streeter-Phelps模型进行水体预测。
        报告结论必须明确总量控制指标是否在许可范围内。
        """

        report = pipeline.expert_training(
            doc_content=ENV_TRAINING_DOC,
            doc_id="eia_report_v2",
            recursion_depth=args.depth
        )

        # Step 3: 知识库搜索
        print("\n\n>>> Step 3: 知识库搜索测试")
        print("─" * 70)
        pipeline.search_knowledge("污染物")

        # Step 4: 文档朗读测试
        print("\n\n>>> Step 4: 朗读测试（读取系统文件）")
        print("─" * 70)

        test_file = r"C:\bak\opencode+omo.txt"
        if os.path.exists(test_file):
            pipeline.read_aloud_file(test_file, max_chars=2000)
        else:
            print(f"  [WARN] 测试文件不存在: {test_file}")
            print(f"  [INFO] 尝试其他测试...")

            # 尝试读取当前目录的 README
            readme = os.path.join(PROJECT_ROOT, "README.md")
            if os.path.exists(readme):
                pipeline.read_aloud_file(readme, max_chars=1000)

        # Step 5: Agent Chat 汇总
        print("\n\n>>> Step 5: Agent Chat 汇总报告")
        print("─" * 70)
        summary_response = pipeline.chat(
            "请总结一下刚才的环评报告专家训练结果，包括知识空白发现数量、"
            "填充情况和知识图谱规模。"
        )
        print(f"\n  汇总响应:\n  {summary_response[:500]}...")

    # 清理
    pipeline.cleanup()
    print("\n\n" + "=" * 70)
    print("  ✅ ExpertTrainingPipeline v2 执行完成")
    print("=" * 70)


if __name__ == "__main__":
    main()
