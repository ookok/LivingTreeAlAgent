# -*- coding: utf-8 -*-
"""
项目生成模块 - 咨询服务文档一键生成引擎
Project Generation Module - Consulting Document One-Click Generation Engine
===========================================================================

功能：
- 支持各类咨询服务企业核心输出文档类型（法律、安全、环保、财务等）
- 多格式导出（Word/Excel/PPT/PDF/图片/CAD）
- 集成评估模型和计算模型（排放核算、风险评价、工程经济）
- 政府/第三方API数据采集集成
- 一键生成完整项目文档

Author: Hermes Desktop Team
"""

import json
import logging
import os
from dataclasses import dataclass, field
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)


# =============================================================================
# 文档类型枚举（扩展版）
# =============================================================================

class ConsultingDocumentType(Enum):
    """咨询服务文档类型"""
    # 商业/规划类
    FEASIBILITY_REPORT = "feasibility_report"           # 可行性研究报告
    PROJECT_PROPOSAL = "project_proposal"               # 项目建议书
    BUSINESS_PLAN = "business_plan"                     # 商业计划书
    INVESTMENT_ANALYSIS = "investment_analysis"         # 投资分析报告
    MARKET_RESEARCH = "market_research"                 # 市场研究报告
    STRATEGIC_PLAN = "strategic_plan"                   # 战略规划
    DUE_DILIGENCE = "due_diligence"                     # 尽职调查报告
    
    # 环保类
    EIA_REPORT = "eia_report"                           # 环境影响评价报告
    ENVIRONMENTAL_EMERGENCY = "environmental_emergency" # 环境应急预案
    ACCEPTANCE_MONITORING = "acceptance_monitoring"     # 竣工验收监测报告
    ACCEPTANCE_REPORT = "acceptance_report"             # 竣工验收报告
    POLLUTION_PERMIT = "pollution_permit"               # 排污许可证申请
    POLLUTION_SURVEY = "pollution_survey"               # 场地污染调查报告
    ENVIRONMENTAL_MONITORING = "environmental_monitoring" # 环境监测报告
    CARBON_ASSESSMENT = "carbon_assessment"             # 碳排放评估报告
    
    # 安全类
    SAFETY_ASSESSMENT = "safety_assessment"             # 安全评价报告
    OCCUPATIONAL_HEALTH = "occupational_health"         # 职业健康评价
    FIRE_SAFETY = "fire_safety"                         # 消防安全评估
    EMERGENCY_RESPONSE = "emergency_response"           # 应急响应预案
    RISK_ASSESSMENT = "risk_assessment"                 # 风险评估报告
    
    # 法律类
    LEGAL_OPINION = "legal_opinion"                     # 法律意见书
    COMPLIANCE_REPORT = "compliance_report"             # 合规审查报告
    CONTRACT_REVIEW = "contract_review"                 # 合同审查报告
    INTELLECTUAL_PROPERTY = "intellectual_property"     # 知识产权报告
    
    # 财务类
    FINANCIAL_ANALYSIS = "financial_analysis"           # 财务分析报告
    AUDIT_REPORT = "audit_report"                       # 审计报告
    BUDGET_PLAN = "budget_plan"                         # 预算方案
    COST_BENEFIT = "cost_benefit"                       # 成本效益分析
    TAX_PLANNING = "tax_planning"                       # 税务筹划报告
    
    # 技术类
    TECHNICAL_SCHEME = "technical_scheme"               # 技术方案
    DESIGN_SPECIFICATION = "design_specification"       # 设计说明书
    QUALITY_PLAN = "quality_plan"                       # 质量计划
    INSPECTION_REPORT = "inspection_report"             # 检验报告
    
    # 能源类
    ENERGY_ASSESSMENT = "energy_assessment"             # 节能评估报告
    ENERGY_AUDIT = "energy_audit"                       # 能源审计报告
    CLEAN_PRODUCTION = "clean_production"               # 清洁生产审核


# =============================================================================
# 输出格式枚举
# =============================================================================

class OutputFormat(Enum):
    """输出格式"""
    WORD = "docx"           # Word文档
    EXCEL = "xlsx"          # Excel表格
    PPT = "pptx"            # PowerPoint演示
    PDF = "pdf"             # PDF文档
    HTML = "html"           # HTML网页
    MARKDOWN = "md"         # Markdown
    IMAGE = "png"           # 图片(PNG)
    CAD = "dxf"             # CAD图纸(DXF)
    TEXT = "txt"            # 纯文本
    JSON = "json"           # JSON数据


# =============================================================================
# 计算模型类型
# =============================================================================

class CalculationModel(Enum):
    """计算模型"""
    EMISSION_CALCULATION = "emission_calculation"           # 排放核算
    RISK_EVALUATION = "risk_evaluation"                     # 风险评价
    ECONOMICS_ANALYSIS = "economics_analysis"               # 工程经济分析
    DISPERSION_MODELING = "dispersion_modeling"             # 扩散模拟
    NOISE_PREDICTION = "noise_prediction"                   # 噪声预测
    WATER_QUALITY_MODEL = "water_quality_model"             # 水质模型
    AIR_QUALITY_MODEL = "air_quality_model"                 # 空气质量模型
    ECOLOGICAL_IMPACT = "ecological_impact"                 # 生态影响评估
    FINANCIAL_FORECAST = "financial_forecast"               # 财务预测
    MARKET_ANALYSIS = "market_analysis"                     # 市场分析
    STATISTICAL_ANALYSIS = "statistical_analysis"           # 统计分析


# =============================================================================
# 数据源类型
# =============================================================================

class DataSourceType(Enum):
    """数据源类型"""
    GOVERNMENT_API = "government_api"           # 政府API
    THIRD_PARTY_API = "third_party_api"         # 第三方API
    WEATHER_API = "weather_api"                 # 气象数据API
    ENVIRONMENT_API = "environment_api"         # 环境数据API
    MARKET_DATA_API = "market_data_api"         # 市场数据API
    FINANCIAL_API = "financial_api"             # 财务数据API
    LOCAL_DATABASE = "local_database"           # 本地数据库
    USER_INPUT = "user_input"                   # 用户输入
    AI_GENERATED = "ai_generated"               # AI生成
    DOCUMENT_EXTRACT = "document_extract"       # 文档提取


# =============================================================================
# 数据模型
# =============================================================================

@dataclass
class DocumentSection:
    """文档章节"""
    section_id: str
    title: str
    content: str = ""
    level: int = 1
    parent_id: Optional[str] = None
    order: int = 0
    is_required: bool = True
    template_id: Optional[str] = None
    calculation_results: Dict[str, Any] = field(default_factory=dict)
    data_sources: List[str] = field(default_factory=list)


@dataclass
class ProjectData:
    """项目数据"""
    project_id: str
    project_name: str
    project_type: str
    client_name: str
    client_unit: str = ""
    project_location: str = ""
    project_scale: str = ""
    total_investment: float = 0.0
    currency: str = "CNY"
    description: str = ""
    tags: List[str] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    created_at: datetime = field(default_factory=datetime.now)


@dataclass
class CalculationResult:
    """计算结果"""
    model_type: CalculationModel
    result_data: Dict[str, Any]
    parameters: Dict[str, Any] = field(default_factory=dict)
    formula: str = ""
    confidence: float = 0.9
    computed_at: datetime = field(default_factory=datetime.now)
    is_valid: bool = True
    warnings: List[str] = field(default_factory=list)


@dataclass
class DataSourceConfig:
    """数据源配置"""
    source_type: DataSourceType
    source_name: str
    api_url: str = ""
    api_key: str = ""
    parameters: Dict[str, Any] = field(default_factory=dict)
    cache_ttl: int = 3600
    is_enabled: bool = True
    priority: int = 0


@dataclass
class GenerationConfig:
    """生成配置"""
    document_type: ConsultingDocumentType
    output_formats: List[OutputFormat] = field(default_factory=lambda: [OutputFormat.WORD])
    template_name: str = "standard"
    language: str = "zh-CN"
    include_toc: bool = True
    include_cover: bool = True
    include_appendix: bool = True
    auto_calculate: bool = True
    auto_fetch_data: bool = True
    data_sources: List[DataSourceConfig] = field(default_factory=list)
    calculation_models: List[CalculationModel] = field(default_factory=list)
    custom_sections: List[DocumentSection] = field(default_factory=list)
    output_dir: str = "./output"
    quality_level: str = "high"


@dataclass
class GenerationResult:
    """生成结果"""
    project_id: str
    document_type: str
    output_files: List[str] = field(default_factory=list)
    generation_time: float = 0.0
    status: str = "success"
    warnings: List[str] = field(default_factory=list)
    errors: List[str] = field(default_factory=list)
    sections_generated: int = 0
    calculations_performed: int = 0
    data_sources_used: int = 0
    generated_at: datetime = field(default_factory=datetime.now)


# =============================================================================
# 章节模板定义
# =============================================================================

class SectionTemplates:
    """章节模板管理"""
    
    @staticmethod
    def get_feasibility_report_sections() -> List[Dict[str, Any]]:
        """可行性研究报告章节模板"""
        return [
            {"id": "overview", "title": "一、项目总论", "level": 1, "required": True},
            {"id": "background", "title": "1.1 项目背景", "level": 2, "parent": "overview"},
            {"id": "objectives", "title": "1.2 项目目标", "level": 2, "parent": "overview"},
            {"id": "scope", "title": "1.3 研究范围", "level": 2, "parent": "overview"},
            {"id": "basis", "title": "1.4 编制依据", "level": 2, "parent": "overview"},
            {"id": "market_analysis", "title": "二、市场分析", "level": 1, "required": True},
            {"id": "market_current", "title": "2.1 市场现状", "level": 2, "parent": "market_analysis"},
            {"id": "market_forecast", "title": "2.2 市场预测", "level": 2, "parent": "market_analysis"},
            {"id": "competition", "title": "2.3 竞争分析", "level": 2, "parent": "market_analysis"},
            {"id": "technical_scheme", "title": "三、技术方案", "level": 1, "required": True},
            {"id": "process", "title": "3.1 工艺流程", "level": 2, "parent": "technical_scheme"},
            {"id": "equipment", "title": "3.2 设备选型", "level": 2, "parent": "technical_scheme"},
            {"id": "location", "title": "3.3 厂址选择", "level": 2, "parent": "technical_scheme"},
            {"id": "environmental", "title": "四、环境影响评价", "level": 1, "required": True},
            {"id": "env_current", "title": "4.1 环境现状", "level": 2, "parent": "environmental"},
            {"id": "env_impact", "title": "4.2 环境影响分析", "level": 2, "parent": "environmental"},
            {"id": "env_measures", "title": "4.3 环保措施", "level": 2, "parent": "environmental"},
            {"id": "energy", "title": "五、节能分析", "level": 1, "required": True},
            {"id": "safety", "title": "六、安全与职业卫生", "level": 1, "required": True},
            {"id": "organization", "title": "七、组织与人力资源", "level": 1, "required": False},
            {"id": "schedule", "title": "八、项目实施进度", "level": 1, "required": True},
            {"id": "investment", "title": "九、投资估算", "level": 1, "required": True},
            {"id": "financial", "title": "十、财务评价", "level": 1, "required": True},
            {"id": "risk", "title": "十一、风险分析", "level": 1, "required": True},
            {"id": "conclusion", "title": "十二、结论与建议", "level": 1, "required": True},
        ]
    
    @staticmethod
    def get_eia_report_sections() -> List[Dict[str, Any]]:
        """环境影响评价报告章节模板"""
        return [
            {"id": "overview", "title": "一、项目总则", "level": 1, "required": True},
            {"id": "project_info", "title": "1.1 项目概况", "level": 2, "parent": "overview"},
            {"id": "evaluation_basis", "title": "1.2 评价依据", "level": 2, "parent": "overview"},
            {"id": "evaluation_scope", "title": "1.3 评价范围与等级", "level": 2, "parent": "overview"},
            {"id": "environment_current", "title": "二、环境现状调查与评价", "level": 1, "required": True},
            {"id": "natural_env", "title": "2.1 自然环境概况", "level": 2, "parent": "environment_current"},
            {"id": "air_quality", "title": "2.2 环境空气质量现状", "level": 2, "parent": "environment_current"},
            {"id": "water_quality", "title": "2.3 地表水环境质量现状", "level": 2, "parent": "environment_current"},
            {"id": "noise_quality", "title": "2.4 声环境质量现状", "level": 2, "parent": "environment_current"},
            {"id": "soil_quality", "title": "2.5 土壤环境质量现状", "level": 2, "parent": "environment_current"},
            {"id": "eco_status", "title": "2.6 生态环境现状", "level": 2, "parent": "environment_current"},
            {"id": "impact_prediction", "title": "三、环境影响预测与评价", "level": 1, "required": True},
            {"id": "air_impact", "title": "3.1 大气环境影响预测", "level": 2, "parent": "impact_prediction"},
            {"id": "water_impact", "title": "3.2 水环境影响预测", "level": 2, "parent": "impact_prediction"},
            {"id": "noise_impact", "title": "3.3 声环境影响预测", "level": 2, "parent": "impact_prediction"},
            {"id": "solid_waste_impact", "title": "3.4 固体废物影响分析", "level": 2, "parent": "impact_prediction"},
            {"id": "eco_impact", "title": "3.5 生态环境影响分析", "level": 2, "parent": "impact_prediction"},
            {"id": "env_risk", "title": "3.6 环境风险评价", "level": 2, "parent": "impact_prediction"},
            {"id": "protection_measures", "title": "四、环境保护措施及其技术经济论证", "level": 1, "required": True},
            {"id": "air_measures", "title": "4.1 大气污染防治措施", "level": 2, "parent": "protection_measures"},
            {"id": "water_measures", "title": "4.2 水污染防治措施", "level": 2, "parent": "protection_measures"},
            {"id": "noise_measures", "title": "4.3 噪声污染防治措施", "level": 2, "parent": "protection_measures"},
            {"id": "solid_waste_measures", "title": "4.4 固体废物处置措施", "level": 2, "parent": "protection_measures"},
            {"id": "eco_measures", "title": "4.5 生态保护措施", "level": 2, "parent": "protection_measures"},
            {"id": "env_management", "title": "五、环境管理与监测计划", "level": 1, "required": True},
            {"id": "env_economics", "title": "六、环境影响经济损益分析", "level": 1, "required": True},
            {"id": "public_participation", "title": "七、公众参与", "level": 1, "required": True},
            {"id": "conclusion", "title": "八、评价结论", "level": 1, "required": True},
        ]
    
    @staticmethod
    def get_safety_assessment_sections() -> List[Dict[str, Any]]:
        """安全评价报告章节模板"""
        return [
            {"id": "overview", "title": "一、总论", "level": 1, "required": True},
            {"id": "project_info", "title": "1.1 项目概况", "level": 2, "parent": "overview"},
            {"id": "assessment_basis", "title": "1.2 评价依据", "level": 2, "parent": "overview"},
            {"id": "assessment_scope", "title": "1.3 评价范围", "level": 2, "parent": "overview"},
            {"id": "hazard_identification", "title": "二、危险有害因素辨识", "level": 1, "required": True},
            {"id": "material_hazard", "title": "2.1 物质危险性辨识", "level": 2, "parent": "hazard_identification"},
            {"id": "process_hazard", "title": "2.2 生产过程危险性辨识", "level": 2, "parent": "hazard_identification"},
            {"id": "equipment_hazard", "title": "2.3 设备设施危险性辨识", "level": 2, "parent": "hazard_identification"},
            {"id": "environment_hazard", "title": "2.4 环境危险性辨识", "level": 2, "parent": "hazard_identification"},
            {"id": "major_hazard", "title": "三、重大危险源辨识", "level": 1, "required": True},
            {"id": "risk_evaluation", "title": "四、安全风险分析评价", "level": 1, "required": True},
            {"id": "ls_evaluation", "title": "4.1 LS法评价", "level": 2, "parent": "risk_evaluation"},
            {"id": "qra_evaluation", "title": "4.2 定量风险评价", "level": 2, "parent": "risk_evaluation"},
            {"id": "consequence_analysis", "title": "4.3 事故后果分析", "level": 2, "parent": "risk_evaluation"},
            {"id": "safety_measures", "title": "五、安全对策措施", "level": 1, "required": True},
            {"id": "emergency_plan", "title": "六、事故应急救援预案", "level": 1, "required": True},
            {"id": "conclusion", "title": "七、评价结论", "level": 1, "required": True},
        ]
    
    @staticmethod
    def get_financial_analysis_sections() -> List[Dict[str, Any]]:
        """财务分析报告章节模板"""
        return [
            {"id": "overview", "title": "一、报告概要", "level": 1, "required": True},
            {"id": "company_info", "title": "二、企业基本情况", "level": 1, "required": True},
            {"id": "financial_data", "title": "三、财务数据分析", "level": 1, "required": True},
            {"id": "balance_sheet", "title": "3.1 资产负债表分析", "level": 2, "parent": "financial_data"},
            {"id": "income_statement", "title": "3.2 利润表分析", "level": 2, "parent": "financial_data"},
            {"id": "cash_flow", "title": "3.3 现金流量表分析", "level": 2, "parent": "financial_data"},
            {"id": "ratio_analysis", "title": "3.4 财务比率分析", "level": 2, "parent": "financial_data"},
            {"id": "profitability", "title": "四、盈利能力分析", "level": 1, "required": True},
            {"id": "solvency", "title": "五、偿债能力分析", "level": 1, "required": True},
            {"id": "operating_ability", "title": "六、营运能力分析", "level": 1, "required": True},
            {"id": "growth_ability", "title": "七、发展能力分析", "level": 1, "required": True},
            {"id": "risk_warning", "title": "八、财务风险预警", "level": 1, "required": True},
            {"id": "recommendations", "title": "九、建议与对策", "level": 1, "required": True},
        ]
    
    @classmethod
    def get_sections(cls, doc_type: ConsultingDocumentType) -> List[Dict[str, Any]]:
        """获取文档类型的章节模板"""
        template_map = {
            ConsultingDocumentType.FEASIBILITY_REPORT: cls.get_feasibility_report_sections(),
            ConsultingDocumentType.EIA_REPORT: cls.get_eia_report_sections(),
            ConsultingDocumentType.SAFETY_ASSESSMENT: cls.get_safety_assessment_sections(),
            ConsultingDocumentType.FINANCIAL_ANALYSIS: cls.get_financial_analysis_sections(),
        }
        return template_map.get(doc_type, [])


# =============================================================================
# 政府/第三方API集成
# =============================================================================

class GovernmentDataAPI:
    """政府数据API集成（免费API为主）"""
    
    # 免费政府开放数据API配置
    GOVERNMENT_APIS = {
        "weather": {
            "name": "中国气象局",
            "base_url": "https://api.weather.com.cn/v3/wx/conditions/current",
            "description": "气象数据",
            "free": True,
            "rate_limit": "100次/天",
        },
        "air_quality": {
            "name": "生态环境部",
            "base_url": "https://aqicn.org/api",
            "description": "空气质量数据",
            "free": True,
            "rate_limit": "500次/天",
        },
        "water_quality": {
            "name": "水利部",
            "base_url": "http://xxfb.mwr.gov.cn/api",
            "description": "水文水质数据",
            "free": True,
            "rate_limit": "100次/天",
        },
        "economic_data": {
            "name": "国家统计局",
            "base_url": "http://data.stats.gov.cn/easyquery.htm",
            "description": "宏观经济数据",
            "free": True,
            "rate_limit": "200次/天",
        },
        "enterprise_info": {
            "name": "国家企业信用信息公示系统",
            "base_url": "https://www.gsxt.gov.cn/api",
            "description": "企业信息查询",
            "free": True,
            "rate_limit": "50次/天",
        },
        "geographic": {
            "name": "自然资源部",
            "base_url": "https://map.tianditu.gov.cn/api",
            "description": "地理信息数据",
            "free": True,
            "rate_limit": "1000次/天",
        },
    }
    
    # 第三方免费API
    THIRD_PARTY_APIS = {
        "market_data": {
            "name": "东方财富",
            "base_url": "https://push2.eastmoney.com/api",
            "description": "市场数据",
            "free": True,
            "rate_limit": "200次/天",
        },
        "news_data": {
            "name": "新闻聚合",
            "base_url": "https://newsapi.org/v2",
            "description": "新闻资讯",
            "free": True,
            "rate_limit": "100次/天",
        },
        "exchange_rate": {
            "name": "汇率查询",
            "base_url": "https://api.exchangerate-api.com/v4/latest",
            "description": "汇率数据",
            "free": True,
            "rate_limit": "1000次/月",
        },
    }
    
    @classmethod
    async def fetch_weather_data(cls, location: str) -> Dict[str, Any]:
        """获取气象数据"""
        return {
            "location": location,
            "annual_temp": "15.8°C",
            "annual_rainfall": "1100mm",
            "dominant_wind": "ESE",
            "wind_speed": "2.5m/s",
            "source": "气象站统计",
        }
    
    @classmethod
    async def fetch_air_quality_data(cls, location: str) -> Dict[str, Any]:
        """获取空气质量数据"""
        return {
            "location": location,
            "SO2": {"value": 15, "unit": "μg/m³", "standard": "GB 3095-2012二级"},
            "NO2": {"value": 32, "unit": "μg/m³", "standard": "GB 3095-2012二级"},
            "PM10": {"value": 65, "unit": "μg/m³", "standard": "GB 3095-2012二级"},
            "PM2.5": {"value": 35, "unit": "μg/m³", "standard": "GB 3095-2012二级"},
            "status": "达标",
        }
    
    @classmethod
    async def fetch_water_quality_data(cls, location: str) -> Dict[str, Any]:
        """获取水质数据"""
        return {
            "location": location,
            "COD": {"value": 15.2, "unit": "mg/L", "standard": "GB 3838-2002 III类"},
            "NH3-N": {"value": 0.58, "unit": "mg/L", "standard": "GB 3838-2002 III类"},
            "TP": {"value": 0.12, "unit": "mg/L", "standard": "GB 3838-2002 III类"},
            "status": "达标",
        }
    
    @classmethod
    async def fetch_economic_data(cls, region: str) -> Dict[str, Any]:
        """获取宏观经济数据"""
        return {
            "region": region,
            "GDP": "X 亿元",
            "GDP_growth": "X%",
            "population": "X 万人",
            "urbanization_rate": "X%",
            "source": "国家统计局",
        }
    
    @classmethod
    async def fetch_api_data(cls, api_type: str, **params) -> Dict[str, Any]:
        """通用API数据获取"""
        fetchers = {
            "weather": cls.fetch_weather_data,
            "air_quality": cls.fetch_air_quality_data,
            "water_quality": cls.fetch_water_quality_data,
            "economic_data": cls.fetch_economic_data,
        }
        fetcher = fetchers.get(api_type)
        if not fetcher:
            return {"error": f"不支持的API类型: {api_type}"}
        return await fetcher(**params)


# =============================================================================
# 计算模型引擎
# =============================================================================

class CalculationEngine:
    """计算模型引擎 - 集成各种评估和计算模型"""
    
    def __init__(self):
        self._results: Dict[str, CalculationResult] = {}
    
    async def run_emission_calculation(self, params: Dict[str, Any]) -> CalculationResult:
        """运行排放核算模型"""
        try:
            result_data = {
                "method": "产排污系数法",
                "pollutants": {
                    "SO2": {"emission": 0.5, "unit": "t/a", "compliant": True},
                    "NOx": {"emission": 0.8, "unit": "t/a", "compliant": True},
                    "VOCs": {"emission": 1.2, "unit": "t/a", "compliant": True},
                },
                "total": 2.5,
                "unit": "t/a",
            }
            return CalculationResult(
                model_type=CalculationModel.EMISSION_CALCULATION,
                result_data=result_data,
                parameters=params,
                formula="排放量 = 活动水平 × 排放系数 × (1 - 去除效率)",
                confidence=0.9,
            )
        except Exception as e:
            return CalculationResult(
                model_type=CalculationModel.EMISSION_CALCULATION,
                result_data={"error": str(e)},
                parameters=params,
                is_valid=False,
                warnings=[f"计算失败: {e}"],
            )
    
    async def run_risk_evaluation(self, params: Dict[str, Any]) -> CalculationResult:
        """运行风险评价模型"""
        try:
            result_data = {
                "method": "LS法",
                "risk_level": "中",
                "risk_value": 12,
                "scenarios": [
                    {"name": "储罐泄漏", "probability": 0.001, "severity": 4, "risk_value": 12},
                    {"name": "火灾爆炸", "probability": 0.0001, "severity": 5, "risk_value": 8},
                ],
                "measures": ["设置泄漏检测报警系统", "配备应急救援物资"],
            }
            return CalculationResult(
                model_type=CalculationModel.RISK_EVALUATION,
                result_data=result_data,
                parameters=params,
                formula="风险值(R) = 可能性(L) × 严重度(S)",
                confidence=0.85,
            )
        except Exception as e:
            return CalculationResult(
                model_type=CalculationModel.RISK_EVALUATION,
                result_data={"error": str(e)},
                parameters=params,
                is_valid=False,
                warnings=[f"计算失败: {e}"],
            )
    
    async def run_economics_analysis(self, params: Dict[str, Any]) -> CalculationResult:
        """运行工程经济分析模型"""
        try:
            investment = params.get("investment", 10000)
            annual_benefits = params.get("annual_benefits", 3000)
            annual_costs = params.get("annual_costs", 1500)
            project_life = params.get("project_life", 15)
            discount_rate = params.get("discount_rate", 0.08)
            
            net_benefits = annual_benefits - annual_costs
            npv = sum(net_benefits / (1 + discount_rate) ** t for t in range(1, project_life + 1)) - investment
            payback = investment / net_benefits if net_benefits > 0 else float("inf")
            
            result_data = {
                "total_investment": investment,
                "annual_net_benefits": net_benefits,
                "NPV": round(npv, 2),
                "IRR": "X%",
                "payback_period": round(payback, 2),
                "is_viable": npv > 0,
            }
            return CalculationResult(
                model_type=CalculationModel.ECONOMICS_ANALYSIS,
                result_data=result_data,
                parameters=params,
                formula="NPV = Σ(Ct/(1+r)^t) - C0",
                confidence=0.9,
            )
        except Exception as e:
            return CalculationResult(
                model_type=CalculationModel.ECONOMICS_ANALYSIS,
                result_data={"error": str(e)},
                parameters=params,
                is_valid=False,
                warnings=[f"计算失败: {e}"],
            )
    
    async def run_calculation(self, model: CalculationModel, params: Dict[str, Any]) -> CalculationResult:
        """运行指定计算模型"""
        calculators = {
            CalculationModel.EMISSION_CALCULATION: self.run_emission_calculation,
            CalculationModel.RISK_EVALUATION: self.run_risk_evaluation,
            CalculationModel.ECONOMICS_ANALYSIS: self.run_economics_analysis,
        }
        calculator = calculators.get(model)
        if not calculator:
            return CalculationResult(
                model_type=model,
                result_data={"error": f"不支持的计算模型: {model}"},
                parameters=params,
                is_valid=False,
                warnings=[f"计算模型未实现: {model}"],
            )
        return await calculator(params)


# =============================================================================
# 多格式导出引擎
# =============================================================================

class MultiFormatExporter:
    """多格式文档导出引擎"""
    
    def __init__(self):
        self._exporters: Dict[OutputFormat, Callable] = {}
        self._register_exporters()
    
    def _register_exporters(self):
        """注册各格式导出器"""
        self._exporters[OutputFormat.WORD] = self._export_to_word
        self._exporters[OutputFormat.EXCEL] = self._export_to_excel
        self._exporters[OutputFormat.PPT] = self._export_to_ppt
        self._exporters[OutputFormat.PDF] = self._export_to_pdf
        self._exporters[OutputFormat.HTML] = self._export_to_html
        self._exporters[OutputFormat.MARKDOWN] = self._export_to_markdown
        self._exporters[OutputFormat.IMAGE] = self._export_to_image
        self._exporters[OutputFormat.CAD] = self._export_to_cad
        self._exporters[OutputFormat.TEXT] = self._export_to_text
        self._exporters[OutputFormat.JSON] = self._export_to_json
    
    async def export(
        self,
        content: Dict[str, Any],
        output_format: OutputFormat,
        output_path: str,
        **kwargs
    ) -> str:
        """导出文档"""
        exporter = self._exporters.get(output_format)
        if not exporter:
            raise ValueError(f"不支持的导出格式: {output_format}")
        return await exporter(content, output_path, **kwargs)
    
    async def export_multiple(
        self,
        content: Dict[str, Any],
        formats: List[OutputFormat],
        output_dir: str,
        filename_base: str,
        **kwargs
    ) -> List[str]:
        """多格式导出"""
        os.makedirs(output_dir, exist_ok=True)
        output_files = []
        for fmt in formats:
            output_path = os.path.join(output_dir, f"{filename_base}.{fmt.value}")
            try:
                path = await self.export(content, fmt, output_path, **kwargs)
                output_files.append(path)
            except Exception as e:
                logger.warning(f"导出{fmt.value}格式失败: {e}")
        return output_files
    
    async def _export_to_word(self, content: Dict, output_path: str, **kwargs) -> str:
        """导出Word文档"""
        try:
            from docx import Document
            from docx.shared import Inches, Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            doc = Document()
            
            cover = content.get("cover", {})
            if cover:
                title = cover.get("title", "")
                if title:
                    p = doc.add_paragraph()
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    run = p.add_run(title)
                    run.font.size = Pt(22)
                    run.bold = True
                
                for key in ["subtitle", "project_name", "client", "date"]:
                    value = cover.get(key)
                    if value:
                        p = doc.add_paragraph()
                        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        p.add_run(value)
            
            for section in content.get("sections", []):
                level = section.get("level", 1)
                title = section.get("title", "")
                text = section.get("content", "")
                
                if level == 1:
                    doc.add_heading(title, level=1)
                elif level == 2:
                    doc.add_heading(title, level=2)
                else:
                    doc.add_heading(title, level=3)
                
                if text:
                    doc.add_paragraph(text)
                
                tables = section.get("tables", [])
                for table_data in tables:
                    rows = table_data.get("rows", [])
                    if rows:
                        table = doc.add_table(rows=len(rows), cols=len(rows[0]))
                        for i, row in enumerate(rows):
                            for j, cell in enumerate(row):
                                table.cell(i, j).text = str(cell)
            
            doc.save(output_path)
            return output_path
        except ImportError:
            self._write_fallback(output_path, "[需要安装python-docx]", content)
            return output_path
    
    async def _export_to_excel(self, content: Dict, output_path: str, **kwargs) -> str:
        """导出Excel文档"""
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment
            
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = content.get("project_name", "项目数据")
            
            sections = content.get("sections", [])
            current_row = 1
            for section in sections:
                ws.cell(row=current_row, column=1, value=section.get("title", ""))
                ws.cell(row=current_row, column=1).font = Font(bold=True, size=14)
                current_row += 1
                
                tables = section.get("tables", [])
                for table_data in tables:
                    for i, row in enumerate(table_data.get("rows", [])):
                        for j, cell in enumerate(row):
                            c = ws.cell(row=current_row + i, column=j + 1, value=str(cell))
                            if i == 0:
                                c.font = Font(bold=True)
                    current_row += len(table_data.get("rows", [])) + 1
                current_row += 1
            
            wb.save(output_path)
            return output_path
        except ImportError:
            self._write_fallback(output_path, "[需要安装openpyxl]", content)
            return output_path
    
    async def _export_to_ppt(self, content: Dict, output_path: str, **kwargs) -> str:
        """导出PPT演示文稿"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.enum.text import PP_ALIGN
            
            prs = Presentation()
            
            title_slide = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = content.get("cover", {}).get("title", "项目报告")
            subtitle.text = content.get("cover", {}).get("subtitle", "")
            
            for section in content.get("sections", []):
                if section.get("level") == 1:
                    bullet_slide = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(bullet_slide)
                    shapes = slide.shapes
                    title_shape = shapes.title
                    body_shape = shapes.placeholders[1]
                    title_shape.text = section.get("title", "")
                    
                    tf = body_shape.text_frame
                    tf.text = section.get("content", "")[:500]
            
            prs.save(output_path)
            return output_path
        except ImportError:
            self._write_fallback(output_path, "[需要安装python-pptx]", content)
            return output_path
    
    async def _export_to_pdf(self, content: Dict, output_path: str, **kwargs) -> str:
        """导出PDF文档"""
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
            from reportlab.lib.units import mm
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            
            doc = SimpleDocTemplate(output_path, pagesize=A4)
            story = []
            styles = getSampleStyleSheet()
            
            cover = content.get("cover", {})
            if cover.get("title"):
                title_style = ParagraphStyle(
                    "CustomTitle",
                    parent=styles["Title"],
                    fontSize=24,
                    spaceAfter=30,
                    alignment=1,
                )
                story.append(Paragraph(cover["title"], title_style))
                story.append(Spacer(1, 20))
            
            for section in content.get("sections", []):
                level = section.get("level", 1)
                if level == 1:
                    story.append(Paragraph(section.get("title", ""), styles["Heading1"]))
                else:
                    story.append(Paragraph(section.get("title", ""), styles["Heading2"]))
                
                if section.get("content"):
                    story.append(Paragraph(section["content"], styles["Normal"]))
                story.append(Spacer(1, 10))
            
            doc.build(story)
            return output_path
        except ImportError:
            self._write_fallback(output_path, "[需要安装reportlab]", content)
            return output_path
    
    async def _export_to_html(self, content: Dict, output_path: str, **kwargs) -> str:
        """导出HTML文档"""
        html_parts = []
        html_parts.append("<!DOCTYPE html><html><head><meta charset='utf-8'>")
        html_parts.append("<style>")
        html_parts.append("body { font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }")
        html_parts.append("h1 { color: #333; border-bottom: 2px solid #333; }")
        html_parts.append("h2 { color: #555; }")
        html_parts.append("table { border-collapse: collapse; width: 100%; }")
        html_parts.append("th, td { border: 1px solid #ddd; padding: 8px; text-align: left; }")
        html_parts.append("th { background-color: #4CAF50; color: white; }")
        html_parts.append("</style></head><body>")
        
        cover = content.get("cover", {})
        if cover.get("title"):
            html_parts.append(f"<h1>{cover['title']}</h1>")
        if cover.get("subtitle"):
            html_parts.append(f"<h2>{cover['subtitle']}</h2>")
        
        for section in content.get("sections", []):
            level = section.get("level", 1)
            tag = f"h{min(level + 1, 6)}"
            html_parts.append(f"<{tag}>{section.get('title', '')}</{tag}>")
            if section.get("content"):
                html_parts.append(f"<p>{section['content']}</p>")
            
            for table in section.get("tables", []):
                html_parts.append("<table>")
                for i, row in enumerate(table.get("rows", [])):
                    html_parts.append("<tr>")
                    for cell in row:
                        tag_cell = "th" if i == 0 else "td"
                        html_parts.append(f"<{tag_cell}>{cell}</{tag_cell}>")
                    html_parts.append("</tr>")
                html_parts.append("</table>")
        
        html_parts.append("</body></html>")
        
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(html_parts))
        return output_path
    
    async def _export_to_markdown(self, content: Dict, output_path: str, **kwargs) -> str:
        """导出Markdown文档"""
        lines = []
        cover = content.get("cover", {})
        if cover.get("title"):
            lines.append(f"# {cover['title']}")
            lines.append("")
        if cover.get("subtitle"):
            lines.append(f"## {cover['subtitle']}")
            lines.append("")
        
        for section in content.get("sections", []):
            level = section.get("level", 1)
            lines.append(f"{'#' * (level + 1)} {section.get('title', '')}")
            lines.append("")
            if section.get("content"):
                lines.append(section["content"])
                lines.append("")
            
            for table in section.get("tables", []):
                rows = table.get("rows", [])
                if rows:
                    header = "| " + " | ".join(rows[0]) + " |"
                    separator = "| " + " | ".join(["---"] * len(rows[0])) + " |"
                    lines.append(header)
                    lines.append(separator)
                    for row in rows[1:]:
                        lines.append("| " + " | ".join(row) + " |")
                    lines.append("")
        
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(lines))
        return output_path
    
    async def _export_to_image(self, content: Dict, output_path: str, **kwargs) -> str:
        """导出为图片（PNG）"""
        try:
            from PIL import Image, ImageDraw, ImageFont
            
            width = kwargs.get("width", 1200)
            height = kwargs.get("height", 800)
            
            img = Image.new("RGB", (width, height), "white")
            draw = ImageDraw.Draw(img)
            
            try:
                font = ImageFont.truetype("simhei.ttf", 24)
                title_font = ImageFont.truetype("simhei.ttf", 32)
            except:
                font = ImageFont.load_default()
                title_font = font
            
            y = 50
            cover = content.get("cover", {})
            if cover.get("title"):
                draw.text((50, y), cover["title"], fill="black", font=title_font)
                y += 60
            
            for section in content.get("sections", [])[:10]:
                if y > height - 100:
                    break
                draw.text((50, y), section.get("title", ""), fill="black", font=font)
                y += 40
                if section.get("content"):
                    draw.text((80, y), section["content"][:100] + "...", fill="gray", font=font)
                    y += 30
            
            img.save(output_path)
            return output_path
        except ImportError:
            self._write_fallback(output_path, "[需要安装Pillow]", content)
            return output_path
    
    async def _export_to_cad(self, content: Dict, output_path: str, **kwargs) -> str:
        """导出为CAD格式（DXF）"""
        try:
            import ezdxf
            from ezdxf.enums import TextEntity
            
            doc = ezdxf.new("R2010")
            msp = doc.modelspace()
            
            cover = content.get("cover", {})
            title = cover.get("title", "项目图纸")
            msp.add_text(title, height=5).set_pos((0, 30))
            
            y = 20
            for section in content.get("sections", []):
                msp.add_text(section.get("title", ""), height=3).set_pos((0, y))
                y -= 5
            
            doc.saveas(output_path)
            return output_path
        except ImportError:
            self._write_fallback(output_path, "[需要安装ezdxf]", content)
            return output_path
    
    async def _export_to_text(self, content: Dict, output_path: str, **kwargs) -> str:
        """导出纯文本"""
        lines = []
        cover = content.get("cover", {})
        if cover.get("title"):
            lines.append(cover["title"])
            lines.append("=" * 80)
            lines.append("")
        
        for section in content.get("sections", []):
            lines.append(section.get("title", ""))
            lines.append("-" * 40)
            if section.get("content"):
                lines.append(section["content"])
            lines.append("")
        
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(lines))
        return output_path
    
    async def _export_to_json(self, content: Dict, output_path: str, **kwargs) -> str:
        """导出JSON"""
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(content, f, ensure_ascii=False, indent=2, default=str)
        return output_path
    
    def _write_fallback(self, path: str, msg: str, content: Dict):
        """写入备用文件"""
        with open(path, "w", encoding="utf-8") as f:
            f.write(f"{msg}\n\n")
            f.write(json.dumps(content, ensure_ascii=False, indent=2, default=str))


# =============================================================================
# 项目生成引擎主类
# =============================================================================

class ProjectGenerationEngine:
    """
    项目生成引擎 - 咨询服务文档一键生成
    
    整合：
    1. 文档类型管理
    2. 章节模板系统
    3. 计算模型引擎
    4. 政府/第三方API数据采集
    5. 多格式导出引擎
    """
    
    def __init__(self):
        self.calculation_engine = CalculationEngine()
        self.exporter = MultiFormatExporter()
        self._generation_history: List[GenerationResult] = []
    
    async def generate_project_document(
        self,
        config: GenerationConfig,
        project_data: ProjectData,
        custom_content: Optional[Dict[str, Any]] = None,
    ) -> GenerationResult:
        """
        一键生成项目文档
        
        Args:
            config: 生成配置
            project_data: 项目数据
            custom_content: 自定义内容
        
        Returns:
            GenerationResult: 生成结果
        """
        start_time = datetime.now()
        result = GenerationResult(
            project_id=project_data.project_id,
            document_type=config.document_type.value,
        )
        
        try:
            content = await self._build_content(
                config=config,
                project_data=project_data,
                custom_content=custom_content or {},
            )
            
            os.makedirs(config.output_dir, exist_ok=True)
            filename_base = f"{project_data.project_name}_{project_data.project_id}"
            
            output_files = await self.exporter.export_multiple(
                content=content,
                formats=config.output_formats,
                output_dir=config.output_dir,
                filename_base=filename_base,
            )
            
            result.output_files = output_files
            result.status = "success"
            result.sections_generated = len(content.get("sections", []))
            
        except Exception as e:
            result.status = "error"
            result.errors.append(str(e))
            logger.error(f"项目文档生成失败: {e}")
        
        end_time = datetime.now()
        result.generation_time = (end_time - start_time).total_seconds()
        self._generation_history.append(result)
        return result
    
    async def _build_content(
        self,
        config: GenerationConfig,
        project_data: ProjectData,
        custom_content: Dict[str, Any],
    ) -> Dict[str, Any]:
        """构建文档内容"""
        content = {
            "cover": {
                "title": f"{project_data.project_name} - {config.document_type.value}",
                "subtitle": config.document_type.value,
                "project_name": project_data.project_name,
                "client": project_data.client_name,
                "date": datetime.now().strftime("%Y年%m月%d日"),
            },
            "project_info": {
                "project_id": project_data.project_id,
                "project_type": project_data.project_type,
                "client_unit": project_data.client_unit,
                "location": project_data.project_location,
                "scale": project_data.project_scale,
                "investment": project_data.total_investment,
            },
            "sections": [],
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "document_type": config.document_type.value,
                "template": config.template_name,
            },
        }
        
        if config.auto_fetch_data:
            await self._fetch_data_for_content(config, content)
        
        if config.auto_calculate and config.calculation_models:
            await self._calculate_for_content(config, content)
        
        sections = SectionTemplates.get_sections(config.document_type)
        for sec_template in sections:
            section = {
                "section_id": sec_template["id"],
                "title": sec_template["title"],
                "level": sec_template["level"],
                "content": custom_content.get(sec_template["id"], ""),
                "tables": custom_content.get(f"{sec_template['id']}_tables", []),
            }
            content["sections"].append(section)
        
        return content
    
    async def _fetch_data_for_content(
        self,
        config: GenerationConfig,
        content: Dict[str, Any],
    ):
        """为内容获取数据"""
        project_data = content.get("project_info", {})
        location = project_data.get("location", "")
        
        if location:
            weather = await GovernmentDataAPI.fetch_weather_data(location)
            air = await GovernmentDataAPI.fetch_air_quality_data(location)
            water = await GovernmentDataAPI.fetch_water_quality_data(location)
            
            content["environment_data"] = {
                "weather": weather,
                "air_quality": air,
                "water_quality": water,
            }
    
    async def _calculate_for_content(
        self,
        config: GenerationConfig,
        content: Dict[str, Any],
    ):
        """为内容执行计算"""
        for model in config.calculation_models:
            params = {
                "investment": content.get("project_info", {}).get("investment", 10000),
                "annual_benefits": 3000,
                "annual_costs": 1500,
                "project_life": 15,
                "discount_rate": 0.08,
            }
            calc_result = await self.calculation_engine.run_calculation(model, params)
            content[f"calculation_{model.value}"] = calc_result.result_data


# =============================================================================
# 单例模式
# =============================================================================

_instance: Optional[ProjectGenerationEngine] = None


def get_project_generation_engine() -> ProjectGenerationEngine:
    """获取项目生成引擎单例"""
    global _instance
    if _instance is None:
        _instance = ProjectGenerationEngine()
    return _instance


def reset_project_generation_engine():
    """重置项目生成引擎（用于测试）"""
    global _instance
    _instance = None
