"""Office API connectors — Google Docs, MS365 Graph, WPS, LaTeX, PPT generation.

Tools exposed to LLM:
  gdocs_create(title, content) → Google Docs document
  ms365_send_email(to, subject, body) → Outlook email
  wps_create(title, content) → WPS document
  export_latex(content) → LaTeX compile to PDF
  export_pptx(title, slides_data) → PowerPoint file
"""

from __future__ import annotations

import os
from pathlib import Path

from loguru import logger


# ═══ P1: Google Docs API ═══

def gdocs_create(title: str, content: str) -> str:
    """Create Google Docs document. Requires GOOGLE_SERVICE_ACCOUNT_JSON env var."""
    creds_file = os.environ.get("GOOGLE_SERVICE_ACCOUNT_JSON", "")
    if not creds_file or not Path(creds_file).exists():
        return (
            "Google Docs not configured. Set GOOGLE_SERVICE_ACCOUNT_JSON to "
            "the path of your service account key file.\n"
            "Get one at: https://console.cloud.google.com/apis/credentials"
        )
    try:
        from google.oauth2 import service_account
        from googleapiclient.discovery import build
        creds = service_account.Credentials.from_service_account_file(
            creds_file, scopes=["https://www.googleapis.com/auth/documents"],
        )
        service = build("docs", "v1", credentials=creds)
        doc = service.documents().create(body={"title": title}).execute()
        doc_id = doc["documentId"]

        # Insert content
        requests = [{"insertText": {"location": {"index": 1}, "text": content}}]
        service.documents().batchUpdate(documentId=doc_id, body={"requests": requests}).execute()

        return f"Google Docs created: https://docs.google.com/document/d/{doc_id}/edit"
    except ImportError:
        return "pip install google-auth google-api-python-client google-auth-httplib2"
    except Exception as e:
        return f"Google Docs error: {e}"


# ═══ P1: MS365 Graph API ═══

def ms365_send_email(to: str, subject: str, body: str) -> str:
    """Send email via MS365 Graph. Requires MS365_CLIENT_ID/TENANT_ID/SECRET env vars."""
    client_id = os.environ.get("MS365_CLIENT_ID", "")
    if not client_id:
        return (
            "MS365 not configured. Set MS365_CLIENT_ID, MS365_TENANT_ID, MS365_CLIENT_SECRET.\n"
            "Register app at: https://portal.azure.com → App registrations"
        )
    try:
        import requests
        tenant = os.environ["MS365_TENANT_ID"]
        secret = os.environ["MS365_CLIENT_SECRET"]
        sender = os.environ.get("MS365_SENDER_EMAIL", "")

        # Get token
        token_resp = requests.post(
            f"https://login.microsoftonline.com/{tenant}/oauth2/v2.0/token",
            data={
                "client_id": client_id, "client_secret": secret,
                "scope": "https://graph.microsoft.com/.default",
                "grant_type": "client_credentials",
            },
            timeout=15,
        )
        token = token_resp.json()["access_token"]

        # Send email
        requests.post(
            "https://graph.microsoft.com/v1.0/me/sendMail",
            headers={"Authorization": f"Bearer {token}", "Content-Type": "application/json"},
            json={
                "message": {
                    "subject": subject,
                    "body": {"contentType": "Text", "content": body},
                    "toRecipients": [{"emailAddress": {"address": addr.strip()}} for addr in to.split(",")],
                },
            },
            timeout=15,
        )
        return f"Email sent to {to} via MS365"
    except ImportError:
        return "pip install requests"
    except Exception as e:
        return f"MS365 error: {e}"


def ms365_create_doc(title: str, content: str) -> str:
    """Create Word doc in OneDrive via MS365 Graph."""
    client_id = os.environ.get("MS365_CLIENT_ID", "")
    if not client_id:
        return "MS365 not configured."
    try:
        import requests
        tenant = os.environ["MS365_TENANT_ID"]
        secret = os.environ["MS365_CLIENT_SECRET"]
        token_resp = requests.post(
            f"https://login.microsoftonline.com/{tenant}/oauth2/v2.0/token",
            data={"client_id": client_id, "client_secret": secret,
                  "scope": "https://graph.microsoft.com/.default",
                  "grant_type": "client_credentials"}, timeout=15,
        )
        token = token_resp.json()["access_token"]
        resp = requests.put(
            f"https://graph.microsoft.com/v1.0/me/drive/root:/{title}.docx:/content",
            headers={"Authorization": f"Bearer {token}"},
            data=content.encode("utf-8"),
            timeout=15,
        )
        return f"MS365 doc created: {title}.docx" if resp.ok else f"MS365 error: {resp.status_code}"
    except Exception as e:
        return f"MS365 error: {e}"


# ═══ P1: WPS API ═══

def wps_create(title: str, content: str) -> str:
    """Create WPS document. Requires WPS_API_KEY env var."""
    api_key = os.environ.get("WPS_API_KEY", "")
    if not api_key:
        return (
            "WPS not configured. Set WPS_API_KEY.\n"
            "Get key at: https://open.wps.cn"
        )
    try:
        import requests
        resp = requests.post(
            "https://openapi.wps.cn/v1/documents",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"title": title, "content": content, "format": "docx"},
            timeout=15,
        )
        if resp.ok:
            data = resp.json()
            return f"WPS doc created: {data.get('url', '')}"
        return f"WPS error: {resp.status_code} {resp.text[:200]}"
    except Exception as e:
        return f"WPS error: {e}"


# ═══ P1: LaTeX Export ═══

def export_latex(content: str, output: str = "output.pdf") -> str:
    """Compile LaTeX to PDF. Requires pdflatex on PATH."""
    import tempfile
    import subprocess as _sp

    # Check pdflatex
    try:
        _sp.run(["pdflatex", "--version"], capture_output=True, timeout=5)
    except FileNotFoundError:
        return "LaTeX not available. Install: sudo apt install texlive-latex-base"

    # Write .tex file
    tex_path = Path(tempfile.mktemp(suffix=".tex"))
    tex_path.write_text(content, encoding="utf-8")

    try:
        result = _sp.run(
            ["pdflatex", "-interaction=nonstopmode", "-output-directory",
             str(tex_path.parent), str(tex_path.name)],
            capture_output=True, text=True, timeout=30, cwd=str(tex_path.parent),
        )
        pdf = tex_path.with_suffix(".pdf")
        if pdf.exists():
            # Move to working dir
            dest = Path(output)
            pdf.rename(dest)
            return f"LaTeX compiled: {dest} ({dest.stat().st_size} bytes)"
        return f"LaTeX compilation failed:\n{result.stdout[-1000:]}"
    except _sp.TimeoutExpired:
        return "LaTeX compilation timed out"
    except Exception as e:
        return f"LaTeX error: {e}"
    finally:
        tex_path.unlink(missing_ok=True)


# ═══ P2: PPT Generation ═══

# ═══ Atomic Typesetting: LLM-controlled docx formatting ═══

def format_docx(spec_json: str) -> str:
    """Generate a precisely-formatted .docx from a JSON spec.

    LLM controls every element atomically — font, size, color, alignment, spacing.

    Spec format:
    {
      "title": "Report Title",
      "page": {"size": "A4", "margin_top_cm": 3.7, "margin_bottom_cm": 3.5,
               "margin_left_cm": 2.8, "margin_right_cm": 2.6},
      "header": {"text": "Company Name — Confidential", "font_size": 9, "align": "center"},
      "footer": {"text": "Page {page}", "font_size": 9, "align": "center"},
      "watermark": {"text": "DRAFT", "font_size": 72, "opacity": 0.1},
      "toc": true,
      "sections": [
        {
          "heading": "1. Introduction",
          "level": 1,
          "content": [
            {"type": "paragraph", "text": "...", "font": "SimSun", "size": 12,
             "bold": false, "italic": false, "color": "#333333",
             "align": "justify", "line_spacing": 1.5, "first_line_indent_cm": 0.74},
            {"type": "table", "headers": ["Col1","Col2"], "rows": [["a","b"]]},
            {"type": "image", "path": "chart.png", "width_cm": 14, "align": "center"}
          ]
        }
      ],
      "output": "report.docx"
    }
    """
    try:
        import json
        from docx import Document
        from docx.shared import Cm, Inches, Pt, RGBColor, Emu
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        spec = json.loads(spec_json)
        doc = Document()
        output = spec.get("output", "output.docx")

        # Page setup
        if "page" in spec:
            pg = spec["page"]
            for section in doc.sections:
                if "size" in pg:
                    section.page_width = Cm({"A4": 21.0, "A3": 29.7, "Letter": 21.59}.get(pg["size"], 21.0))
                    section.page_height = Cm({"A4": 29.7, "A3": 42.0, "Letter": 27.94}.get(pg["size"], 29.7))
                section.top_margin = Cm(pg.get("margin_top_cm", 2.54))
                section.bottom_margin = Cm(pg.get("margin_bottom_cm", 2.54))
                section.left_margin = Cm(pg.get("margin_left_cm", 3.18))
                section.right_margin = Cm(pg.get("margin_right_cm", 3.18))

        # Header/Footer
        for section in doc.sections:
            if "header" in spec:
                hdr = spec["header"]
                header = section.header
                header.is_linked_to_previous = False
                p = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
                p.text = hdr.get("text", "")
                p.alignment = _parse_align(hdr.get("align", "center"))
                if hdr.get("font_size"):
                    for run in p.runs:
                        run.font.size = Pt(hdr["font_size"])

            if "footer" in spec:
                ftr = spec["footer"]
                footer = section.footer
                footer.is_linked_to_previous = False
                p = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
                _add_page_number(p)
                p.alignment = _parse_align(ftr.get("align", "center"))
                if ftr.get("font_size"):
                    for run in p.runs:
                        run.font.size = Pt(ftr["font_size"])

        # Title
        if "title" in spec:
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run(spec["title"])
            run.font.size = Pt(22)
            run.font.bold = True

        # TOC
        if spec.get("toc"):
            doc.add_paragraph("Table of Contents", style="Heading 1")
            for sec in spec.get("sections", []):
                p = doc.add_paragraph()
                p.add_run(f"{sec.get('heading','')}").font.size = Pt(11)

        # Sections
        for sec in spec.get("sections", []):
            heading_text = sec.get("heading", "")
            level = sec.get("level", 1)
            if heading_text:
                doc.add_heading(heading_text, level=min(level, 3))

            for item in sec.get("content", []):
                item_type = item.get("type", "paragraph")

                if item_type == "paragraph":
                    p = doc.add_paragraph()
                    _apply_paragraph_format(p, item)

                elif item_type == "table":
                    headers = item.get("headers", [])
                    rows = item.get("rows", [])
                    if headers or rows:
                        table = doc.add_table(rows=len(rows) + 1, cols=max(len(headers), 1))
                        table.style = "Table Grid"
                        for i, h in enumerate(headers):
                            cell = table.rows[0].cells[i]
                            cell.text = h
                            for run in cell.paragraphs[0].runs:
                                run.font.bold = True
                        for r, row in enumerate(rows):
                            for c, val in enumerate(row[:len(headers)]):
                                table.rows[r+1].cells[c].text = str(val)

                elif item_type == "image":
                    img_path = item.get("path", "")
                    if img_path and Path(img_path).exists():
                        width = item.get("width_cm", 14)
                        p = doc.add_paragraph()
                        p.alignment = _parse_align(item.get("align", "center"))
                        p.add_run().add_picture(img_path, width=Cm(width))

        # Watermark
        if "watermark" in spec:
            wm = spec["watermark"]
            for section in doc.sections:
                _add_watermark(section, wm.get("text", ""), wm.get("font_size", 72))

        doc.save(output)
        size = Path(output).stat().st_size
        return f"Formatted docx saved: {output} ({size} bytes, {len(spec.get('sections',[]))} sections)"
    except ImportError:
        return "pip install python-docx"
    except Exception as e:
        return f"format_docx error: {e}"


def _parse_align(align: str):
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    return {"left": 0, "center": 1, "right": 2, "justify": 3}.get(align, 1)


def _apply_paragraph_format(p, item: dict):
    from docx.shared import Cm, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn

    text = item.get("text", "")
    font_name = item.get("font", "SimSun")
    font_size = item.get("size", 12)
    bold = item.get("bold", False)
    italic = item.get("italic", False)
    color = item.get("color", "#000000")
    align = item.get("align", "justify")
    line_spacing = item.get("line_spacing", 1.5)
    first_indent = item.get("first_line_indent_cm", 0)

    run = p.add_run(text)
    run.font.name = font_name
    run._element.rPr.rFonts.set(qn("w:eastAsia"), font_name)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color.startswith("#"):
        run.font.color.rgb = RGBColor.from_string(color[1:])
    p.alignment = _parse_align(align)

    # Line spacing
    pf = p.paragraph_format
    pf.line_spacing = line_spacing
    if first_indent > 0:
        pf.first_line_indent = Cm(first_indent)


def _add_page_number(paragraph):
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    run = paragraph.add_run()
    fld_char_begin = OxmlElement("w:fldChar")
    fld_char_begin.set(qn("w:fldCharType"), "begin")
    run._r.append(fld_char_begin)
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = " PAGE "
    run._r.append(instr)
    fld_char_end = OxmlElement("w:fldChar")
    fld_char_end.set(qn("w:fldCharType"), "end")
    run._r.append(fld_char_end)


def _add_watermark(section, text: str, font_size: int = 72):
    from docx.shared import Pt, RGBColor
    from docx.oxml.ns import qn
    header = section.header
    header.is_linked_to_previous = False
    p = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
    p.alignment = 1
    run = p.add_run(text)
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor(200, 200, 200)

def export_pptx(title: str, slides_json: str) -> str:
    """Generate PowerPoint from JSON slide data.

    slides_json format: [{"title":"Slide 1","content":["bullet1","bullet2"]}, ...]
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        import json

        slides = json.loads(slides_json)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title
        if slides:
            title_slide.placeholders[1].text = slides[0].get("subtitle", "")

        # Content slides
        for i, slide_data in enumerate(slides[:20]):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_data.get("title", f"Slide {i+1}")

            body = slide.placeholders[1].text_frame
            body.clear()
            for bullet in slide_data.get("content", [])[:10]:
                p = body.add_paragraph()
                p.text = str(bullet)
                p.level = 0

        out = Path(title.replace(" ", "_") + ".pptx")
        prs.save(str(out))
        return f"PPT created: {out} ({len(slides)} slides, {out.stat().st_size} bytes)"
    except ImportError:
        return "pip install python-pptx"
    except Exception as e:
        return f"PPT error: {e}"


# ── Tool registration ──

TOOLS = {
    "gdocs_create": {"desc": "Create Google Docs document.", "params": "title, content"},
    "ms365_send_email": {"desc": "Send email via MS365 Outlook.", "params": "to, subject, body"},
    "ms365_create_doc": {"desc": "Create Word doc in OneDrive via MS365.", "params": "title, content"},
    "wps_create": {"desc": "Create WPS document.", "params": "title, content"},
    "export_latex": {"desc": "Compile LaTeX to PDF.", "params": "latex_content [output_path]"},
    "format_docx": {"desc": "Generate precisely-formatted .docx with atomic control over fonts, margins, headers, watermarks, TOC.", "params": "json_spec"}, "export_pptx": {"desc": "Generate PowerPoint from JSON slides.", "params": "title, slides_json"},
}
