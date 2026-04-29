"""
office_preview/office_renderer.py - Office 文档渲染引擎

支持 Word (.docx) / Excel (.xlsx, .csv) / PowerPoint (.pptx) 预览
采用 HTML 渲染方式，兼容 PyQt6 QWebEngineView
"""

import os
import io
import csv
from typing import Optional, List, Dict, Any, Tuple
from dataclasses import dataclass
from .models import RenderResult, WordPage, ExcelSheet, PowerPointSlide


class WordRenderer:
    """Word 文档渲染器"""

    def __init__(self):
        self._docx_available = self._check_docx()

    def _check_docx(self) -> bool:
        try:
            import docx
            return True
        except ImportError:
            return False

    def render(self, file_path: str) -> RenderResult:
        """渲染 Word 文档为 HTML"""
        if not os.path.exists(file_path):
            return RenderResult.error(f"文件不存在: {file_path}")

        if not self._docx_available:
            return self._render_plain_text_fallback(file_path)

        try:
            import docx
            from docx.shared import Pt, RGBColor, Inches
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            doc = docx.Document(file_path)
            pages = self._extract_pages(doc)
            html = self._build_word_html(doc, pages)
            metadata = {
                'title': doc.core_properties.title or '',
                'author': doc.core_properties.author or '',
                'page_count': len(pages),
                'word_count': sum(p.word_count for p in pages),
            }
            return RenderResult.ok(html, metadata)

        except Exception as e:
            return self._render_plain_text_fallback(file_path)

    def _extract_pages(self, doc) -> List[WordPage]:
        """提取文档页面信息（基于段落分段）"""
        pages = []
        current_page = WordPage(
            page_number=1,
            heading='',
            content_preview='',
            word_count=0
        )
        word_count = 0

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            word_count += len(text.split())

            # 检测标题
            if para.style.name.startswith('Heading'):
                if current_page.heading:
                    current_page.word_count = word_count
                    pages.append(current_page)
                    current_page = WordPage(
                        page_number=len(pages) + 1,
                        heading=text,
                        content_preview='',
                        word_count=0
                    )
                    word_count = 0
                else:
                    current_page.heading = text
                    current_page.content_preview = text

        if current_page.heading or word_count > 0:
            current_page.word_count = word_count
            pages.append(current_page)

        return pages if pages else [WordPage(1, '文档', '', 0)]

    def _build_word_html(self, doc, pages: List[WordPage]) -> str:
        """构建 Word HTML"""
        # 提取所有段落内容
        paragraphs_html = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                paragraphs_html.append('<br/>')
                continue

            style = para.style.name
            if style.startswith('Heading 1'):
                paragraphs_html.append(f'<h1>{self._escape_html(text)}</h1>')
            elif style.startswith('Heading 2'):
                paragraphs_html.append(f'<h2>{self._escape_html(text)}</h2>')
            elif style.startswith('Heading 3'):
                paragraphs_html.append(f'<h3>{self._escape_html(text)}</h3>')
            elif style == 'Title':
                paragraphs_html.append(f'<h1 class="doc-title">{self._escape_html(text)}</h1>')
            elif style == 'Quote':
                paragraphs_html.append(f'<blockquote><p>{self._escape_html(text)}</p></blockquote>')
            else:
                # 处理段落内的富文本
                runs_html = []
                for run in para.runs:
                    r_text = run.text
                    if r_text.strip():
                        if run.bold and run.italic:
                            runs_html.append(
                                f'<strong><em>{self._escape_html(r_text)}</em></strong>'
                            )
                        elif run.bold:
                            runs_html.append(f'<strong>{self._escape_html(r_text)}</strong>')
                        elif run.italic:
                            runs_html.append(f'<em>{self._escape_html(r_text)}</em>')
                        else:
                            runs_html.append(self._escape_html(r_text))
                paragraphs_html.append(f'<p>{"".join(runs_html)}</p>')

        # 提取表格
        tables_html = []
        for i, table in enumerate(doc.tables):
            tables_html.append(self._render_table(table))

        return f'''<!DOCTYPE html>
<html><head><meta charset="utf-8">
<style>
:root {{
    --bg: #ffffff;
    --fg: #1a1a1a;
    --accent: #2563eb;
    --border: #e5e7eb;
    --code-bg: #f3f4f6;
}}
* {{ box-sizing: border-box; margin: 0; padding: 0; }}
body {{ font-family: 'Segoe UI', -apple-system, sans-serif; font-size: 14px;
        line-height: 1.8; color: var(--fg); padding: 32px;
        max-width: 900px; margin: 0 auto; background: var(--bg); }}
h1 {{ font-size: 1.8em; color: #111; margin: 1em 0 0.5em; border-bottom: 2px solid var(--accent); padding-bottom: 0.3em; }}
h2 {{ font-size: 1.4em; color: #333; margin: 0.8em 0 0.4em; }}
h3 {{ font-size: 1.15em; color: #555; }}
.doc-title {{ text-align: center; font-size: 2em; color: #111; margin-bottom: 1em; }}
p {{ margin: 0.6em 0; }}
blockquote {{ border-left: 4px solid var(--accent); padding: 4px 16px; margin: 1em 0;
             background: #f8fafc; }}
table {{ border-collapse: collapse; width: 100%; margin: 1em 0; font-size: 13px; }}
th, td {{ border: 1px solid var(--border); padding: 8px 12px; text-align: left; }}
th {{ background: #f8fafc; font-weight: 600; }}
tr:nth-child(even) {{ background: #f9fafb; }}
.page-nav {{ position: fixed; bottom: 20px; right: 20px; display: flex; gap: 8px;
              background: #fff; border: 1px solid var(--border); border-radius: 8px;
              padding: 8px 12px; box-shadow: 0 2px 8px rgba(0,0,0,0.1); }}
.page-nav button {{ background: var(--accent); color: #fff; border: none; padding: 6px 12px;
                    border-radius: 4px; cursor: pointer; }}
</style></head><body>
<div class="doc-content">
{chr(10).join(paragraphs_html)}
{chr(10).join(tables_html)}
</div>
</body></html>'''

    def _render_table(self, table) -> str:
        """渲染 Word 表格"""
        rows = []
        for i, row in enumerate(table.rows):
            cells = [cell.text.strip() for cell in row.cells]
            tag = 'th' if i == 0 else 'td'
            rows.append(f'<tr>{"".join(f"<{tag}>{self._escape_html(c)}</{tag}>" for c in cells)}</tr>')
        return f'<table><thead>{"".join(rows[:1])}</thead><tbody>{"".join(rows[1:])}</tbody></table>'

    def _escape_html(self, text: str) -> str:
        return (text.replace('&', '&amp;').replace('<', '&lt;')
                .replace('>', '&gt;').replace('"', '&quot;'))

    def _render_plain_text_fallback(self, file_path: str) -> RenderResult:
        """纯文本降级渲染"""
        try:
            with open(file_path, 'rb') as f:
                raw = f.read()
            # 尝试提取文本
            text = self._extract_text_from_docx_raw(raw)
            html = f'''<!DOCTYPE html><html><head><meta charset="utf-8">
<style>body {{ font-family: monospace; padding: 24px; white-space: pre-wrap;
            word-break: break-all; }}</style></head><body>{self._escape_html(text)}</body></html>'''
            return RenderResult.ok(html, {'fallback': True})
        except Exception as e:
            return RenderResult.error(f"无法渲染 Word 文档: {str(e)}")

    def _extract_text_from_docx_raw(self, raw: bytes) -> str:
        """从 DOCX 二进制流中提取文本"""
        import zipfile
        try:
            with zipfile.ZipFile(io.BytesIO(raw)) as z:
                with z.open('word/document.xml') as f:
                    import re
                    xml = f.read().decode('utf-8')
                    text = re.sub(r'<[^>]+>', ' ', xml)
                    text = ' '.join(text.split())
                    return text[:5000]
        except Exception:
            return "无法提取文档内容"


class ExcelRenderer:
    """Excel 文档渲染器"""

    def __init__(self):
        self._openpyxl_available = self._check_openpyxl()

    def _check_openpyxl(self) -> bool:
        try:
            import openpyxl
            return True
        except ImportError:
            return False

    def render(self, file_path: str) -> RenderResult:
        """渲染 Excel 文档为 HTML"""
        ext = os.path.splitext(file_path)[1].lower()

        if ext == '.csv':
            return self._render_csv(file_path)

        if not self._openpyxl_available:
            return self._render_csv_fallback(file_path)

        try:
            import openpyxl
            from openpyxl.utils import get_column_letter

            wb = openpyxl.load_workbook(file_path, data_only=True)
            sheets = self._extract_sheets(wb)
            html = self._build_excel_html(wb, sheets)
            metadata = {
                'sheet_count': len(sheets),
                'sheet_names': [s.name for s in sheets],
                'total_rows': sum(s.row_count for s in sheets),
            }
            return RenderResult.ok(html, metadata)

        except Exception as e:
            return self._render_csv_fallback(file_path)

    def _extract_sheets(self, wb) -> List[ExcelSheet]:
        """提取工作表信息"""
        sheets = []
        for name in wb.sheetnames[:20]:  # 最多20个表
            ws = wb[name]
            row_count = ws.max_row or 0
            col_count = ws.max_column or 0

            # 预览数据（前10行）
            preview = []
            for row in ws.iter_rows(min_row=1, max_row=10, max_col=min(col_count, 10), values_only=True):
                preview.append([str(c) if c is not None else '' for c in row])

            sheets.append(ExcelSheet(
                name=name,
                row_count=row_count,
                column_count=col_count,
                used_range=f"A1:{get_column_letter(col_count)}{row_count}" if col_count and row_count else "",
                preview_data=preview
            ))
        return sheets

    def _build_excel_html(self, wb, sheets: List[ExcelSheet]) -> str:
        """构建 Excel HTML"""
        # Sheet 选择器
        tabs_html = []
        for i, s in enumerate(sheets):
            active = 'active' if i == 0 else ''
            tabs_html.append(
                f'<button class="sheet-tab {active}" data-sheet="{i}">'
                f'{self._escape_html(s.name)}</button>'
            )

        # 内容
        sheets_html = []
        for i, s in enumerate(sheets):
            display = 'block' if i == 0 else 'none'
            rows_html = []

            # 表头
            if s.preview_data:
                header = s.preview_data[0]
                rows_html.append('<thead><tr>')
                for h in header:
                    rows_html.append(f'<th>{self._escape_html(h)}</th>')
                rows_html.append('</tr></thead><tbody>')

                # 数据行
                for row in s.preview_data[1:]:
                    rows_html.append('<tr>')
                    for cell in row:
                        rows_html.append(f'<td>{self._escape_html(cell)}</td>')
                    rows_html.append('</tr>')
                rows_html.append('</tbody>')

            sheets_html.append(
                f'<div class="sheet-content" id="sheet-{i}" style="display:{display}">'
                f'<table>{chr(10).join(rows_html)}</table>'
                f'<p class="sheet-info">共 {s.row_count} 行 × {s.column_count} 列 | '
                f'范围: {s.used_range}</p>'
                f'</div>'
            )

        return f'''<!DOCTYPE html>
<html><head><meta charset="utf-8">
<style>
:root {{ --border: #d0d7de; --header-bg: #f6f8fa; --row-alt: #f6f8fa; }}
* {{ box-sizing: border-box; margin: 0; padding: 0; }}
body {{ font-family: 'Segoe UI', -apple-system, sans-serif; font-size: 13px;
        background: #fff; color: #1a1a1a; }}
.sheet-tabs {{ display: flex; gap: 4px; padding: 12px 16px; background: var(--header-bg);
               border-bottom: 1px solid var(--border); overflow-x: auto; }}
.sheet-tab {{ background: #eaeef2; border: 1px solid var(--border); border-radius: 6px;
              padding: 6px 16px; cursor: pointer; font-size: 13px; white-space: nowrap;
              border-bottom: none; border-radius: 6px 6px 0 0; }}
.sheet-tab.active {{ background: #fff; border-bottom: 2px solid #fff;
                     margin-bottom: -1px; font-weight: 600; }}
.sheet-content {{ padding: 16px; }}
table {{ border-collapse: collapse; width: 100%; font-size: 13px; }}
th, td {{ border: 1px solid var(--border); padding: 6px 10px; text-align: left;
          max-width: 300px; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; }}
th {{ background: var(--header-bg); font-weight: 600; position: sticky; top: 0; }}
tr:nth-child(even) td {{ background: var(--row-alt); }}
tr:hover td {{ background: #e8f4fd; }}
.sheet-info {{ margin-top: 12px; font-size: 12px; color: #666; }}
</style>
</head><body>
<div class="sheet-tabs">{"".join(tabs_html)}</div>
<div class="sheets">{"".join(sheets_html)}</div>
<script>
document.querySelectorAll('.sheet-tab').forEach(btn => {{
    btn.addEventListener('click', () => {{
        document.querySelectorAll('.sheet-tab').forEach(b => b.classList.remove('active'));
        document.querySelectorAll('.sheet-content').forEach(s => s.style.display = 'none');
        btn.classList.add('active');
        document.getElementById('sheet-' + btn.dataset.sheet).style.display = 'block';
    }});
}});
</script>
</body></html>'''

    def _render_csv(self, file_path: str) -> RenderResult:
        """渲染 CSV 文件"""
        try:
            with open(file_path, 'r', encoding='utf-8-sig', errors='replace') as f:
                reader = csv.reader(f)
                rows = list(reader)[:100]  # 最多100行

            if not rows:
                return RenderResult.ok('<html><body><p>空文件</p></body></html>')

            header = rows[0] if rows else []
            data_rows = rows[1:]

            html_rows = ['<thead><tr>']
            for h in header:
                html_rows.append(f'<th>{self._escape_html(h)}</th>')
            html_rows.append('</tr></thead><tbody>')

            for row in data_rows:
                html_rows.append('<tr>')
                for cell in row:
                    html_rows.append(f'<td>{self._escape_html(cell)}</td>')
                html_rows.append('</tr>')
            html_rows.append('</tbody>')

            html = f'''<!DOCTYPE html><html><head><meta charset="utf-8">
<style>
* {{ box-sizing: border-box; }} body {{ font-family: 'Segoe UI', sans-serif;
  font-size: 13px; }} table {{ border-collapse: collapse; width: 100%; }}
th, td {{ border: 1px solid #d0d7de; padding: 6px 10px; text-align: left; }}
th {{ background: #f6f8fa; font-weight: 600; }} tr:nth-child(even) td
{{ background: #f6f8fa; }}
</style></head><body>
<table>{chr(10).join(html_rows)}</table>
<p style="padding:12px;color:#666;font-size:12px;">显示 {len(data_rows)} 行，共 {len(rows)} 行</p>
</body></html>'''
            return RenderResult.ok(html, {'rows': len(rows), 'columns': len(header)})

        except Exception as e:
            return RenderResult.error(f"CSV 渲染失败: {str(e)}")

    def _render_csv_fallback(self, file_path: str) -> RenderResult:
        """CSV 降级渲染"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read()[:5000]
            return RenderResult.ok(
                f'<html><body><pre style="font-family:monospace;padding:16px;">'
                f'{self._escape_html(content)}</pre></body></html>',
                {'fallback': True}
            )
        except Exception as e:
            return RenderResult.error(f"文件读取失败: {str(e)}")

    def _escape_html(self, text: str) -> str:
        return (str(text).replace('&', '&amp;').replace('<', '&lt;')
                .replace('>', '&gt;').replace('"', '&quot;'))


class PowerPointRenderer:
    """PowerPoint 渲染器"""

    def __init__(self):
        self._pptx_available = self._check_pptx()

    def _check_pptx(self) -> bool:
        try:
            import pptx
            return True
        except ImportError:
            return False

    def render(self, file_path: str) -> RenderResult:
        """渲染 PowerPoint 为 HTML"""
        if not os.path.exists(file_path):
            return RenderResult.error(f"文件不存在: {file_path}")

        if not self._pptx_available:
            return self._render_fallback(file_path)

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            prs = Presentation(file_path)
            slides = self._extract_slides(prs)
            html = self._build_slides_html(prs, slides)
            metadata = {
                'slide_count': len(slides),
                'slide_titles': [s.title for s in slides if s.title],
            }
            return RenderResult.ok(html, metadata)

        except Exception as e:
            return self._render_fallback(file_path)

    def _extract_slides(self, prs) -> List[PowerPointSlide]:
        """提取幻灯片信息"""
        slides = []
        for i, slide in enumerate(prs.slides[:50], 1):  # 最多50张
            title = ''
            content = []

            # 标题
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()

            # 内容
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text and text != title:
                            content.append(text)
                            if len(content) >= 10:
                                break

            slides.append(PowerPointSlide(
                slide_number=i,
                title=title or f'第 {i} 张幻灯片',
                content_preview=' | '.join(content[:5]),
                image_count=sum(1 for s in slide.shapes if s.shape_type == 13),
                layout=slide.slide_layout.name if slide.slide_layout else ''
            ))
        return slides

    def _build_slides_html(self, prs, slides: List[PowerPointSlide]) -> str:
        """构建幻灯片 HTML"""
        # 幻灯片缩略图区域
        thumbs_html = []
        content_html = []

        for i, slide in enumerate(slides):
            # 缩略图
            thumbs_html.append(
                f'<div class="slide-thumb" data-slide="{i}">'
                f'<div class="thumb-number">{slide.slide_number}</div>'
                f'<div class="thumb-title">{self._escape_html(slide.title)}</div>'
                f'</div>'
            )

            # 内容
            content_html.append(
                f'<div class="slide-content" id="slide-content-{i}">'
                f'<div class="slide-header">'
                f'<span class="slide-num">幻灯片 {slide.slide_number} / {len(slides)}</span>'
                f'<span class="slide-layout">{self._escape_html(slide.layout)}</span>'
                f'</div>'
                f'<h2 class="slide-title">{self._escape_html(slide.title)}</h2>'
                f'<p class="slide-preview">{self._escape_html(slide.content_preview)}</p>'
                f'</div>'
            )

        return f'''<!DOCTYPE html>
<html><head><meta charset="utf-8">
<style>
:root {{ --bg: #1a1a2e; --panel-bg: #16213e; --accent: #0f3460; --text: #eaeaea;
         --highlight: #e94560; --border: #333; }}
* {{ box-sizing: border-box; margin: 0; padding: 0; }}
body {{ font-family: 'Segoe UI', -apple-system, sans-serif; background: var(--bg);
        color: var(--text); display: flex; height: 100vh; overflow: hidden; }}
.slide-list {{ width: 220px; background: var(--panel-bg); border-right: 1px solid var(--border);
               overflow-y: auto; padding: 12px; }}
.slide-list h3 {{ font-size: 13px; color: #888; margin-bottom: 12px; text-transform: uppercase;
                  letter-spacing: 1px; }}
.slide-thumb {{ background: var(--accent); border-radius: 8px; padding: 12px; margin-bottom: 8px;
                cursor: pointer; transition: all 0.2s; border: 2px solid transparent; }}
.slide-thumb:hover {{ border-color: var(--highlight); transform: translateX(4px); }}
.slide-thumb.active {{ border-color: var(--highlight); background: var(--highlight); }}
.thumb-number {{ font-size: 11px; opacity: 0.7; }}
.thumb-title {{ font-size: 13px; font-weight: 500; margin-top: 4px; overflow: hidden;
                text-overflow: ellipsis; white-space: nowrap; }}
.slide-viewer {{ flex: 1; display: flex; flex-direction: column; overflow: hidden; }}
.slide-content {{ flex: 1; padding: 32px; overflow-y: auto; display: none; }}
.slide-content.active {{ display: block; }}
.slide-header {{ display: flex; justify-content: space-between; font-size: 12px; color: #888;
                 margin-bottom: 16px; }}
.slide-title {{ font-size: 1.8em; color: #fff; margin-bottom: 16px; border-bottom: 2px solid var(--highlight);
                padding-bottom: 12px; }}
.slide-preview {{ font-size: 15px; line-height: 1.8; color: #ccc; white-space: pre-wrap; }}
</style>
</head><body>
<div class="slide-list">
    <h3>幻灯片 ({len(slides)}张)</h3>
    {''.join(thumbs_html)}
</div>
<div class="slide-viewer">
    {''.join(content_html)}
</div>
<script>
document.querySelectorAll('.slide-thumb').forEach(thumb => {{
    thumb.addEventListener('click', () => {{
        document.querySelectorAll('.slide-thumb').forEach(t => t.classList.remove('active'));
        document.querySelectorAll('.slide-content').forEach(c => c.classList.remove('active'));
        thumb.classList.add('active');
        document.getElementById('slide-content-' + thumb.dataset.slide).classList.add('active');
    }});
}});
document.querySelector('.slide-thumb').click();
</script>
</body></html>'''

    def _render_fallback(self, file_path: str) -> RenderResult:
        """降级渲染"""
        return RenderResult.ok(
            f'<html><body style="padding:32px;font-family:monospace;">'
            f'<p>PowerPoint 文件: {self._escape_html(os.path.basename(file_path))}</p>'
            f'<p>需要 python-pptx 库才能预览内容</p>'
            f'<p>安装: pip install python-pptx</p></body></html>',
            {'fallback': True}
        )

    def _escape_html(self, text: str) -> str:
        return (str(text).replace('&', '&amp;').replace('<', '&lt;')
                .replace('>', '&gt;').replace('"', '&quot;'))


class OfficeRenderer:
    """Office 文档统一渲染器"""

    def __init__(self):
        self.word = WordRenderer()
        self.excel = ExcelRenderer()
        self.powerpoint = PowerPointRenderer()

    def render(self, file_path: str) -> RenderResult:
        """根据文件类型选择渲染器"""
        ext = os.path.splitext(file_path)[1].lower()

        if ext in {'.docx', '.doc', '.odt'}:
            return self.word.render(file_path)
        elif ext in {'.xlsx', '.xls', '.xlsm', '.ods', '.csv'}:
            return self.excel.render(file_path)
        elif ext in {'.pptx', '.ppt', '.odp'}:
            return self.powerpoint.render(file_path)
        else:
            return RenderResult.error(f"不支持的 Office 格式: {ext}")
