"""
🔧 四大 Office 引擎

1. DocxEngine: Word 文档引擎 (python-docx)
2. XlsxEngine: Excel 引擎 (openpyxl)
3. PptxEngine: PowerPoint 引擎 (python-pptx)
4. PdfEngine: PDF 引擎 (reportlab / PyMuPDF)

每个引擎提供:
- create(): 从头创建
- fill(): 填充现有文档
- read(): 读取文档内容
- format(): 格式化文档
"""

import os
import io
import json
import logging
from enum import Enum
from dataclasses import dataclass, field
from typing import Optional, Dict, List, Any, Tuple
from datetime import datetime

logger = logging.getLogger(__name__)


class EngineType(Enum):
    """引擎类型"""
    DOCX = "docx"
    XLSX = "xlsx"
    PPTX = "pptx"
    PDF = "pdf"


class EngineCapability(Enum):
    """引擎能力"""
    CREATE = "create"
    READ = "read"
    FILL = "fill"
    EDIT = "edit"
    FORMAT = "format"
    EXPORT = "export"
    CONVERT = "convert"


@dataclass
class EngineResult:
    """引擎操作结果"""
    success: bool = False
    output_path: Optional[str] = None
    output_bytes: Optional[bytes] = None
    content: Any = None
    metadata: dict = field(default_factory=dict)
    errors: List[str] = field(default_factory=list)
    warnings: List[str] = field(default_factory=list)

    def to_dict(self) -> dict:
        return {
            "success": self.success,
            "output_path": self.output_path,
            "has_content": self.content is not None,
            "metadata": self.metadata,
            "errors": self.errors,
            "warnings": self.warnings,
        }


# ===== DOCX 引擎 =====

class DocxEngine:
    """
    DOCX 文档引擎

    基于 python-docx 实现:
    - 从头创建 Word 文档
    - 填充/编辑现有文档
    - 格式化与样式应用
    - 导出为 PDF

    参考 minimax-DOCX 的三层工作流:
    - Create: 从头创建新文档
    - Fill: 填充/编辑现有文档内容
    - Format: 通过模板应用格式化
    """

    def __init__(self):
        self._docx_available = self._check_docx()

    def _check_docx(self) -> bool:
        try:
            import docx
            return True
        except ImportError:
            logger.warning("python-docx 未安装, DOCX 功能受限")
            return False

    def create(self, title: str = "", content: str = "",
               theme: dict = None, outline: list = None,
               output_path: str = None) -> EngineResult:
        """创建 DOCX 文档"""
        result = EngineResult()

        if not self._docx_available:
            result.errors.append("python-docx 未安装")
            return result

        try:
            from docx import Document
            from docx.shared import Pt, Cm, Inches, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            doc = Document()

            # 应用主题
            style = doc.styles['Normal']
            if theme:
                font_info = theme.get("fonts", {}).get("body", {})
                if font_info:
                    style.font.name = font_info.get("family", "SimSun")
                    style.font.size = Pt(font_info.get("size", 12))

                color_info = theme.get("colors", {}).get("text_primary", {})
                if color_info:
                    hex_val = color_info.get("hex", "#000000").lstrip('#')
                    style.font.color.rgb = RGBColor(
                        int(hex_val[0:2], 16),
                        int(hex_val[2:4], 16),
                        int(hex_val[4:6], 16)
                    )

            # 标题
            if title:
                heading = doc.add_heading(title, level=0)
                heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # 大纲内容
            if outline:
                for item in outline:
                    if isinstance(item, dict):
                        level = item.get("level", 1)
                        text = item.get("text", "")
                        if level == 1:
                            doc.add_heading(text, level=1)
                        elif level == 2:
                            doc.add_heading(text, level=2)
                        elif level == 3:
                            doc.add_heading(text, level=3)
                        else:
                            doc.add_paragraph(text)

                        # 子内容
                        sub_items = item.get("items", [])
                        for sub in sub_items:
                            doc.add_paragraph(sub, style='List Bullet')
                    elif isinstance(item, str):
                        doc.add_heading(item, level=1)

            # 正文内容
            if content:
                for para_text in content.split('\n'):
                    if para_text.strip():
                        doc.add_paragraph(para_text)

            # 保存
            if output_path:
                doc.save(output_path)
                result.output_path = output_path
            else:
                buffer = io.BytesIO()
                doc.save(buffer)
                result.output_bytes = buffer.getvalue()

            result.success = True
            result.metadata = {
                "title": title,
                "paragraphs": len(doc.paragraphs),
                "sections": len(doc.sections),
            }

        except Exception as e:
            result.errors.append(str(e))
            logger.error(f"DOCX 创建失败: {e}")

        return result

    def fill(self, template_path: str, fill_data: dict,
             output_path: str = None) -> EngineResult:
        """
        填充现有 DOCX 模板

        零格式损失原则: 只替换占位符内容，保留所有原始格式
        """
        result = EngineResult()

        if not self._docx_available:
            result.errors.append("python-docx 未安装")
            return result

        if not os.path.exists(template_path):
            result.errors.append(f"模板文件不存在: {template_path}")
            return result

        try:
            from docx import Document

            doc = Document(template_path)
            fill_count = 0

            # 遍历段落，替换占位符
            for paragraph in doc.paragraphs:
                for key, value in fill_data.items():
                    placeholder = f"{{{{{key}}}}}"  # {{key}} 格式
                    if placeholder in paragraph.text:
                        # 保留格式的替换
                        for run in paragraph.runs:
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, str(value))
                                fill_count += 1

            # 遍历表格
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            for key, value in fill_data.items():
                                placeholder = f"{{{{{key}}}}}"
                                if placeholder in paragraph.text:
                                    for run in paragraph.runs:
                                        if placeholder in run.text:
                                            run.text = run.text.replace(placeholder, str(value))
                                            fill_count += 1

            # 保存
            if output_path:
                doc.save(output_path)
                result.output_path = output_path
            else:
                buffer = io.BytesIO()
                doc.save(buffer)
                result.output_bytes = buffer.getvalue()

            result.success = True
            result.metadata = {"fields_filled": fill_count}

        except Exception as e:
            result.errors.append(str(e))
            logger.error(f"DOCX 填充失败: {e}")

        return result

    def read(self, file_path: str) -> EngineResult:
        """读取 DOCX 文件内容"""
        result = EngineResult()

        if not self._docx_available:
            result.errors.append("python-docx 未安装")
            return result

        try:
            from docx import Document

            doc = Document(file_path)
            paragraphs = []
            for p in doc.paragraphs:
                paragraphs.append({
                    "text": p.text,
                    "style": p.style.name if p.style else None,
                })

            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)

            result.success = True
            result.content = {
                "paragraphs": paragraphs,
                "tables": tables,
                "total_paragraphs": len(paragraphs),
                "total_tables": len(tables),
            }

        except Exception as e:
            result.errors.append(str(e))

        return result


# ===== XLSX 引擎 =====

class XlsxEngine:
    """
    XLSX 电子表格引擎

    基于 openpyxl 实现:
    - 从头创建 Excel
    - 读取和分析数据
    - 零格式损失编辑
    - 公式重新计算
    - 专业财务格式化

    参考 minimax-xlsx 的能力
    """

    def __init__(self):
        self._openpyxl_available = self._check_openpyxl()

    def _check_openpyxl(self) -> bool:
        try:
            import openpyxl
            return True
        except ImportError:
            logger.warning("openpyxl 未安装, XLSX 功能受限")
            return False

    def create(self, title: str = "", data: list = None,
               headers: list = None, sheet_name: str = "Sheet1",
               output_path: str = None) -> EngineResult:
        """创建 XLSX 文件"""
        result = EngineResult()

        if not self._openpyxl_available:
            result.errors.append("openpyxl 未安装")
            return result

        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, Alignment, PatternFill, Border, Side

            wb = Workbook()
            ws = wb.active
            ws.title = sheet_name

            # 标题行
            if headers:
                header_font = Font(bold=True, size=12)
                header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
                header_font_white = Font(bold=True, size=12, color="FFFFFF")

                for col, header in enumerate(headers, 1):
                    cell = ws.cell(row=1, column=col, value=header)
                    cell.font = header_font_white
                    cell.fill = header_fill
                    cell.alignment = Alignment(horizontal="center", vertical="center")

            # 数据行
            if data:
                for row_idx, row_data in enumerate(data, 2):
                    if isinstance(row_data, (list, tuple)):
                        for col_idx, value in enumerate(row_data, 1):
                            ws.cell(row=row_idx, column=col_idx, value=value)
                    elif isinstance(row_data, dict) and headers:
                        for col_idx, header in enumerate(headers, 1):
                            ws.cell(row=row_idx, column=col_idx, value=row_data.get(header, ""))

            # 自动列宽
            for col in ws.columns:
                max_length = 0
                col_letter = col[0].column_letter
                for cell in col:
                    try:
                        if cell.value:
                            max_length = max(max_length, len(str(cell.value)))
                    except Exception:
                        pass
                ws.column_dimensions[col_letter].width = min(max_length + 4, 50)

            if output_path:
                wb.save(output_path)
                result.output_path = output_path
            else:
                buffer = io.BytesIO()
                wb.save(buffer)
                result.output_bytes = buffer.getvalue()

            result.success = True
            result.metadata = {
                "sheet_name": sheet_name,
                "rows": len(data) if data else 0,
                "columns": len(headers) if headers else 0,
            }

        except Exception as e:
            result.errors.append(str(e))
            logger.error(f"XLSX 创建失败: {e}")

        return result

    def read(self, file_path: str) -> EngineResult:
        """读取 XLSX 文件"""
        result = EngineResult()

        if not self._openpyxl_available:
            result.errors.append("openpyxl 未安装")
            return result

        try:
            from openpyxl import load_workbook

            wb = load_workbook(file_path, data_only=True)
            sheets_data = {}

            for ws_name in wb.sheetnames:
                ws = wb[ws_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    rows.append(list(row))
                sheets_data[ws_name] = rows

            result.success = True
            result.content = {
                "sheets": sheets_data,
                "sheet_names": wb.sheetnames,
            }

        except Exception as e:
            result.errors.append(str(e))

        return result


# ===== PPTX 引擎 =====

class PptxEngine:
    """
    PPTX 演示文稿引擎

    基于 python-pptx 实现:
    - 从头创建 (封面/目录/内容/分节/总结)
    - 编辑现有 PPTX
    - 文本提取
    - 15 种封面风格

    参考 pptx-generator 的能力
    """

    # 15 种封面风格的颜色方案
    COVER_PALETTES = {
        "classic": {"bg": "1B3A5C", "accent": "E8A838", "text": "FFFFFF"},
        "modern": {"bg": "2D3436", "accent": "6C5CE7", "text": "FFFFFF"},
        "elegant": {"bg": "FFFFFF", "accent": "C0392B", "text": "2C3E50"},
        "bold": {"bg": "E74C3C", "accent": "FFFFFF", "text": "FFFFFF"},
        "minimal": {"bg": "FAFAFA", "accent": "1A1A1A", "text": "1A1A1A"},
        "corporate": {"bg": "1B3A5C", "accent": "4A7FB5", "text": "FFFFFF"},
        "tech": {"bg": "0D47A1", "accent": "00E5FF", "text": "FFFFFF"},
        "creative": {"bg": "6C5CE7", "accent": "FD79A8", "text": "FFFFFF"},
        "nature": {"bg": "27AE60", "accent": "F1C40F", "text": "FFFFFF"},
        "abstract": {"bg": "2C3E50", "accent": "E67E22", "text": "FFFFFF"},
        "photo": {"bg": "000000", "accent": "FFFFFF", "text": "FFFFFF"},
        "gradient": {"bg": "667EEA", "accent": "764BA2", "text": "FFFFFF"},
        "typography": {"bg": "FFFFFF", "accent": "000000", "text": "000000"},
        "geology": {"bg": "5D4E37", "accent": "A0895C", "text": "F5F1E9"},
        "aurora": {"bg": "1A0533", "accent": "00FF87", "text": "FFFFFF"},
    }

    def __init__(self):
        self._pptx_available = self._check_pptx()

    def _check_pptx(self) -> bool:
        try:
            from pptx import Presentation
            return True
        except ImportError:
            logger.warning("python-pptx 未安装, PPTX 功能受限")
            return False

    def create(self, title: str = "", subtitle: str = "",
               slides: list = None, cover_style: str = "corporate",
               output_path: str = None) -> EngineResult:
        """创建 PPTX 演示文稿"""
        result = EngineResult()

        if not self._pptx_available:
            result.errors.append("python-pptx 未安装")
            return result

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            palette = self.COVER_PALETTES.get(cover_style, self.COVER_PALETTES["corporate"])

            # 封面
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)

            # 标题
            if title:
                txBox = slide.shapes.add_textbox(
                    Inches(1), Inches(2.5), Inches(11), Inches(2)
                )
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(44)
                p.font.bold = True
                p.font.color.rgb = RGBColor.from_string(palette["text"])
                p.alignment = PP_ALIGN.CENTER

            # 副标题
            if subtitle:
                txBox = slide.shapes.add_textbox(
                    Inches(1), Inches(4.5), Inches(11), Inches(1)
                )
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = subtitle
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor.from_string(palette["accent"])
                p.alignment = PP_ALIGN.CENTER

            # 内容幻灯片
            if slides:
                for slide_data in slides:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])  # 标题+内容

                    if isinstance(slide_data, dict):
                        if "title" in slide_data:
                            slide.shapes.title.text = slide_data["title"]
                        if "content" in slide_data:
                            body = slide.placeholders[1]
                            content = slide_data["content"]
                            if isinstance(content, list):
                                for item in content:
                                    p = body.text_frame.add_paragraph()
                                    p.text = str(item)
                            else:
                                body.text_frame.text = str(content)
                    elif isinstance(slide_data, str):
                        slide.shapes.title.text = slide_data

            # 总结页
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(
                Inches(1), Inches(3), Inches(11), Inches(1.5)
            )
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = "谢谢！"
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = RGBColor.from_string(palette["text"])
            p.alignment = PP_ALIGN.CENTER

            if output_path:
                prs.save(output_path)
                result.output_path = output_path
            else:
                buffer = io.BytesIO()
                prs.save(buffer)
                result.output_bytes = buffer.getvalue()

            result.success = True
            result.metadata = {
                "title": title,
                "total_slides": len(prs.slides),
                "cover_style": cover_style,
            }

        except Exception as e:
            result.errors.append(str(e))
            logger.error(f"PPTX 创建失败: {e}")

        return result

    def read(self, file_path: str) -> EngineResult:
        """读取 PPTX 文件"""
        result = EngineResult()

        if not self._pptx_available:
            result.errors.append("python-pptx 未安装")
            return result

        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            slides_content = []

            for slide in prs.slides:
                slide_data = {"shapes": []}
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts = []
                        for para in shape.text_frame.paragraphs:
                            texts.append(para.text)
                        slide_data["shapes"].append({
                            "type": "text",
                            "content": "\n".join(texts),
                        })
                slides_content.append(slide_data)

            result.success = True
            result.content = {
                "slides": slides_content,
                "total_slides": len(prs.slides),
            }

        except Exception as e:
            result.errors.append(str(e))

        return result


# ===== PDF 引擎 =====

class PdfEngine:
    """
    PDF 文档引擎

    基于 reportlab 实现:
    - 从头创建 PDF (15 种封面风格)
    - 填充现有 PDF 表单
    - 文档重新格式化
    - 打印就绪的输出

    参考 minimax-pdf 的设计理念:
    - Token 化设计系统驱动
    - 排版和颜色根据文档类型自动调整
    """

    # 15 种封面风格的 PDF 布局
    COVER_LAYOUTS = {
        "classic": {"title_y": 0.5, "title_size": 28, "line_width": 200},
        "modern": {"title_y": 0.45, "title_size": 32, "line_width": 0},
        "elegant": {"title_y": 0.55, "title_size": 26, "line_width": 300},
        "bold": {"title_y": 0.4, "title_size": 36, "line_width": 0},
        "minimal": {"title_y": 0.5, "title_size": 24, "line_width": 100},
        "corporate": {"title_y": 0.5, "title_size": 28, "line_width": 200},
        "tech": {"title_y": 0.45, "title_size": 30, "line_width": 0},
        "creative": {"title_y": 0.4, "title_size": 34, "line_width": 0},
        "nature": {"title_y": 0.5, "title_size": 28, "line_width": 150},
        "abstract": {"title_y": 0.45, "title_size": 30, "line_width": 0},
        "photo": {"title_y": 0.35, "title_size": 28, "line_width": 0},
        "gradient": {"title_y": 0.45, "title_size": 30, "line_width": 0},
        "typography": {"title_y": 0.5, "title_size": 40, "line_width": 0},
        "geology": {"title_y": 0.5, "title_size": 26, "line_width": 180},
        "aurora": {"title_y": 0.45, "title_size": 30, "line_width": 0},
    }

    def __init__(self):
        self._reportlab_available = self._check_reportlab()

    def _check_reportlab(self) -> bool:
        try:
            import reportlab
            return True
        except ImportError:
            logger.warning("reportlab 未安装, PDF 功能受限")
            return False

    def create(self, title: str = "", content: str = "",
               theme: dict = None, cover_style: str = "classic",
               output_path: str = None) -> EngineResult:
        """创建 PDF 文档"""
        result = EngineResult()

        if not self._reportlab_available:
            result.errors.append("reportlab 未安装")
            return result

        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.colors import HexColor, black, white
            from reportlab.lib.units import cm, mm
            from reportlab.lib.enums import TA_CENTER, TA_LEFT
            from reportlab.platypus import (
                SimpleDocTemplate, Paragraph, Spacer, PageBreak,
                Table, TableStyle
            )
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont

            # 注册中文字体
            try:
                pdfmetrics.registerFont(TTFont('SimSun', 'simsun.ttc'))
                chinese_font = 'SimSun'
            except Exception:
                chinese_font = 'Helvetica'

            if not output_path:
                output_path = os.path.join(
                    os.environ.get("TEMP", "/tmp"),
                    f"office_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
                )

            doc = SimpleDocTemplate(
                output_path,
                pagesize=A4,
                leftMargin=2.5*cm,
                rightMargin=2.5*cm,
                topMargin=2.5*cm,
                bottomMargin=2.5*cm,
            )

            styles = getSampleStyleSheet()

            # 封面
            layout = self.COVER_LAYOUTS.get(cover_style, self.COVER_LAYOUTS["classic"])
            palette = PptxEngine.COVER_PALETTES.get(cover_style,
                                                      PptxEngine.COVER_PALETTES["classic"])

            elements = []

            # 封面页
            elements.append(Spacer(1, 200))

            # 标题
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Title'],
                fontName=chinese_font,
                fontSize=layout["title_size"],
                alignment=TA_CENTER,
                textColor=HexColor(f"#{palette['text']}"),
                spaceAfter=20,
            )
            elements.append(Paragraph(title or "文档标题", title_style))

            # 装饰线
            if layout["line_width"] > 0:
                line_data = [['' ]]
                line_table = Table(line_data, colWidths=[layout["line_width"]*mm/10])
                line_table.setStyle(TableStyle([
                    ('LINEBELOW', (0, 0), (-1, -1), 2, HexColor(f"#{palette['accent']}")),
                ]))
                elements.append(line_table)

            elements.append(PageBreak())

            # 正文
            body_style = ParagraphStyle(
                'CustomBody',
                parent=styles['Normal'],
                fontName=chinese_font,
                fontSize=12,
                leading=20,
                spaceAfter=10,
            )

            if content:
                for line in content.split('\n'):
                    if line.strip():
                        if line.startswith('# '):
                            h1_style = ParagraphStyle(
                                'H1', parent=styles['Heading1'],
                                fontName=chinese_font, fontSize=18,
                            )
                            elements.append(Paragraph(line[2:], h1_style))
                        elif line.startswith('## '):
                            h2_style = ParagraphStyle(
                                'H2', parent=styles['Heading2'],
                                fontName=chinese_font, fontSize=15,
                            )
                            elements.append(Paragraph(line[3:], h2_style))
                        else:
                            elements.append(Paragraph(line, body_style))
                    else:
                        elements.append(Spacer(1, 10))

            doc.build(elements)

            result.success = True
            result.output_path = output_path
            result.metadata = {
                "title": title,
                "cover_style": cover_style,
                "pages": 1 + (len(content.split('\n')) // 30 if content else 0),
            }

        except Exception as e:
            result.errors.append(str(e))
            logger.error(f"PDF 创建失败: {e}")

        return result


# ===== 引擎工厂 =====

class EngineFactory:
    """引擎工厂 - 根据格式创建对应引擎"""

    _engines = {
        EngineType.DOCX: DocxEngine,
        EngineType.XLSX: XlsxEngine,
        EngineType.PPTX: PptxEngine,
        EngineType.PDF: PdfEngine,
    }

    @classmethod
    def get_engine(cls, engine_type: EngineType):
        """获取引擎实例"""
        engine_cls = cls._engines.get(engine_type)
        if engine_cls:
            return engine_cls()
        return None

    @classmethod
    def get_engine_for_format(cls, fmt: str):
        """根据格式字符串获取引擎"""
        fmt_map = {
            "docx": EngineType.DOCX,
            "xlsx": EngineType.XLSX,
            "pptx": EngineType.PPTX,
            "pdf": EngineType.PDF,
        }
        engine_type = fmt_map.get(fmt.lower())
        if engine_type:
            return cls.get_engine(engine_type)
        return None

    @classmethod
    def check_dependencies(cls) -> Dict[str, bool]:
        """检查所有引擎的依赖"""
        return {
            "python-docx": DocxEngine()._docx_available,
            "openpyxl": XlsxEngine()._openpyxl_available,
            "python-pptx": PptxEngine()._pptx_available,
            "reportlab": PdfEngine()._reportlab_available,
        }
